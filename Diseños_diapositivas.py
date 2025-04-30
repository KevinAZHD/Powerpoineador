import os, sys
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageEnhance, ImageFilter

# Clase para manejar los diseños de diapositivas
class Diapositivas:
    def __init__(self, presentation, title_font_name='Calibri', content_font_name='Calibri', title_font_size=16, content_font_size=10, title_bold=False, title_italic=False, title_underline=False, content_bold=False, content_italic=False, content_underline=False):
        self.presentation = presentation
        self.title_font_name = title_font_name
        self.content_font_name = content_font_name
        self.title_font_size = title_font_size
        self.content_font_size = content_font_size
        self.title_bold = title_bold
        self.title_italic = title_italic
        self.title_underline = title_underline
        self.content_bold = content_bold
        self.content_italic = content_italic
        self.content_underline = content_underline

    # Función auxiliar para aplicar fuente al texto
    def apply_font(self, text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.title_font_name

    # Nuevas funciones para aplicar fuentes específicas
    def apply_title_font(self, text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.title_font_name

    def apply_content_font(self, text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.content_font_name
                # Aplicar formato de contenido
                if self.content_bold: run.font.bold = True
                else: run.font.bold = False # Asegurar que no se quede activo si no está marcado
                if self.content_italic: run.font.italic = True
                else: run.font.italic = False
                if self.content_underline: run.font.underline = True
                else: run.font.underline = False

    # Diseño de introducción - Moderno y profesional
    def design0(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Procesar imagen para fondo con efecto de desenfoque
        try:
            if sys.platform == 'win32':
                APP_DATA_DIR = os.path.join(os.getenv('APPDATA'), 'Powerpoineador')
            elif sys.platform == 'darwin':
                APP_DATA_DIR = os.path.join(os.path.expanduser('~'), 'Library', 'Application Support', 'Powerpoineador')
            else:
                APP_DATA_DIR = os.path.join(os.path.expanduser('~'), '.Powerpoineador')
            
            IMAGES_DIR = os.path.join(APP_DATA_DIR, 'images')
            
            if not os.path.exists(IMAGES_DIR):
                os.makedirs(IMAGES_DIR, exist_ok=True)
                
            image = Image.open(image_path)
            blurred_image = image.filter(ImageFilter.GaussianBlur(radius=3))
            enhanced = ImageEnhance.Brightness(blurred_image).enhance(0.7)
            processed_path = os.path.join(IMAGES_DIR, 'intro_background.jpg')
            enhanced.save(processed_path)
            bg_image_path = processed_path
        except Exception as e:
            print(f"Error al procesar imagen: {str(e)}")
            bg_image_path = image_path

        # Añadir imagen de fondo
        left = top = Inches(0)
        width = self.presentation.slide_width
        height = self.presentation.slide_height

        # Añadir gradiente usando dos rectángulos con transparencia
        gradient_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height/2)
        fill_top = gradient_top.fill
        fill_top.solid()
        fill_top.fore_color.rgb = RGBColor(242, 101, 34)
        gradient_top.transparency = 0.6
        gradient_top.line.fill.background()

        gradient_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, height/2, width, height/2)
        fill_bottom = gradient_bottom.fill
        fill_bottom.solid()
        fill_bottom.fore_color.rgb = RGBColor(150, 50, 0)
        gradient_bottom.transparency = 0.5
        gradient_bottom.line.fill.background()

        # Añadir barra lateral decorativa
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), height)
        sidebar_fill = sidebar.fill
        sidebar_fill.solid()
        sidebar_fill.fore_color.rgb = RGBColor(242, 101, 34)
        sidebar.line.fill.background()

        # Añadir imagen visible con borde decorativo
        img_side = min(height/3, width/3)
        img_left = width - img_side - Inches(1)
        img_top = Inches(0.6)
        
        # Borde decorativo para la imagen
        img_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left - Inches(0.1),
            img_top - Inches(0.1),
            img_side + Inches(0.2),
            img_side + Inches(0.2)
        )
        border_fill = img_border.fill
        border_fill.solid()
        border_fill.fore_color.rgb = RGBColor(252, 180, 65)
        img_border.line.width = Pt(3)
        img_border.line.color.rgb = RGBColor(255, 255, 255)
        
        # Imagen principal visible
        visible_pic = slide.shapes.add_picture(image_path, img_left, img_top, width=img_side, height=img_side)
        
        # Añadir título con estilo moderno
        title_left = Inches(0.7)
        title_top = Inches(0.3)
        title_width = img_left - title_left - Inches(0.5)
        title_height = Inches(2)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        
        title_frame.auto_size = MSO_AUTO_SIZE.NONE
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.margin_bottom = Inches(0.1)
        title_frame.margin_left = Inches(0.1)
        title_frame.margin_right = Inches(0.1)
        
        title_frame.text = section
        p_title = title_frame.paragraphs[0]
        p_title.alignment = PP_ALIGN.LEFT
        run_title = p_title.runs[0]
        run_title.font.size = Pt(self.title_font_size)
        self.apply_title_font(title_frame)
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(255, 255, 255)

        # Añadir línea decorativa bajo el título
        underline_top_fixed = Inches(3.4)
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            title_left,
            underline_top_fixed,
            min(Inches(6), title_width),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(255, 126, 0)
        title_underline.line.fill.background()

        # Añadir texto de contenido
        text_left = Inches(1.2)
        text_top = Inches(3.6)
        text_width = img_left - text_left - Inches(-2)
        text_height = Inches(3)
        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        
        text_frame.clear()
        p_content = text_frame.add_paragraph()
        p_content.text = content
        p_content.alignment = PP_ALIGN.LEFT
        run_content = p_content.runs[0]
        run_content.font.size = Pt(self.content_font_size)
        run_content.font.color.rgb = RGBColor(240, 240, 240)
        self.apply_content_font(text_frame)

        # Añadir elementos decorativos
        # Forma en esquina inferior derecha
        corner_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            width - Inches(1.2),
            height - Inches(1.2),
            Inches(1),
            Inches(1)
        )
        corner_fill = corner_shape.fill
        corner_fill.solid()
        corner_fill.fore_color.rgb = RGBColor(255, 140, 0)
        corner_shape.line.fill.background()
        corner_shape.rotation = 45

        # Pequeño icono decorativo
        icon = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            height - Inches(1),
            Inches(0.4),
            Inches(0.5)
        )
        icon_fill = icon.fill
        icon_fill.solid()
        icon_fill.fore_color.rgb = RGBColor(252, 180, 65)
        icon.line.fill.background()

    # Diseño que muestra una imagen a la izquierda y un título y contenido a la derecha con estilo moderno
    def design1(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Añadir fondo con color gradiente para aspecto profesional
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), 
            Inches(0), 
            self.presentation.slide_width, 
            self.presentation.slide_height
        )
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Añadir una imagen vertical a la izquierda con efecto de sombra
        left = Inches(0)
        top = Inches(0)
        width = Inches(5)
        height = Inches(7.55)
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        
        # Agregar barra decorativa vertical en la separación entre imagen y contenido
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5),
            Inches(0),
            Inches(0.1),
            self.presentation.slide_height
        )
        bar_fill = bar.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = RGBColor(0, 120, 215)
        bar.line.fill.background()
        
        # Área de contenido con fondo sutilmente destacado
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.1),
            Inches(0),
            self.presentation.slide_width - Inches(5.1),
            self.presentation.slide_height
        )
        content_bg_fill = content_bg.fill
        content_bg_fill.solid()
        content_bg_fill.fore_color.rgb = RGBColor(250, 250, 250)
        content_bg.transparency = 0.2
        content_bg.line.fill.background()

        # Añadir un título a la derecha con estilo moderno
        title_box = slide.shapes.add_textbox(Inches(5.3), Inches(0.2), Inches(3.3), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].runs[0].font.size = Pt(self.title_font_size)
        title_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 120, 215)
        title_frame.margin_bottom = Inches(0)
        title_frame.margin_left = Inches(0.25)
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(0, 120, 215)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.3),
            Inches(2.1),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(0, 120, 215)
        title_underline.line.fill.background()

        # Añadir texto debajo del título con estilo mejorado
        text_box = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4), Inches(4))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(50, 50, 50)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)
        
        # Elemento decorativo en esquina superior derecha
        decoration = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.presentation.slide_width - Inches(1.2),
            Inches(0.3),
            Inches(1),
            Inches(1)
        )
        dec_fill = decoration.fill
        dec_fill.solid()
        dec_fill.fore_color.rgb = RGBColor(0, 120, 215)
        decoration.line.fill.background()
        decoration.rotation = 45

    # Diseño que muestra una imagen a la derecha y un título y contenido a la izquierda con estilo moderno
    def design2(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Añadir fondo con color gradiente para aspecto profesional
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), 
            Inches(0), 
            self.presentation.slide_width, 
            self.presentation.slide_height
        )
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Añadir una imagen vertical a la derecha con efecto de sombra
        left = Inches(5)
        top = Inches(0)
        width = Inches(5)
        height = Inches(7.55)
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        
        # Agregar barra decorativa vertical en la separación entre imagen y contenido
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.9),
            Inches(0),
            Inches(0.1),
            self.presentation.slide_height
        )
        bar_fill = bar.fill
        bar_fill.solid()
        bar_fill.fore_color.rgb = RGBColor(33, 115, 70)
        bar.line.fill.background()
        
        # Área de contenido con fondo sutilmente destacado
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(4.9),
            self.presentation.slide_height
        )
        content_bg_fill = content_bg.fill
        content_bg_fill.solid()
        content_bg_fill.fore_color.rgb = RGBColor(250, 250, 250)
        content_bg.transparency = 0.2
        content_bg.line.fill.background()

        # Añadir un título a la izquierda con estilo moderno
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.2), Inches(3.3), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].runs[0].font.size = Pt(self.title_font_size)
        title_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(33, 115, 70)
        title_frame.margin_bottom = Inches(0)
        title_frame.margin_left = Inches(0.25)
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(33, 115, 70)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4),
            Inches(2.1),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(33, 115, 70)
        title_underline.line.fill.background()

        # Añadir texto debajo del título con estilo mejorado
        text_box = slide.shapes.add_textbox(Inches(0.4), Inches(2), Inches(4), Inches(4))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(50, 50, 50)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)
        
        # Elemento decorativo en esquina superior izquierda
        decoration = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.2),
            Inches(0.3),
            Inches(1),
            Inches(1)
        )
        dec_fill = decoration.fill
        dec_fill.solid()
        dec_fill.fore_color.rgb = RGBColor(33, 115, 70)
        decoration.line.fill.background()
        decoration.rotation = 45

    # Diseño que muestra una imagen para que ocupe toda la diapositiva y un título y contenido a la derecha
    def design3(self, slide, section, content, image_path):
        slide_bg_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_bg_layout)

        # Añadir una imagen para que ocupe toda la diapositiva
        left = Inches(0)
        top = Inches(0)
        width = self.presentation.slide_width
        height = self.presentation.slide_height
        # Identificar la imagen como fondo con un nombre específico para identificación posterior
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        pic.name = "SLIDE_BACKGROUND_IMAGE"  # Nombre especial para identificarla

        # Añadir una forma de relleno con esquinas redondeadas y semitransparente a la derecha
        left = Inches(4.5)
        top = Inches(1.2)
        width = Inches(5)
        height = Inches(5.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 250)
        shape.line.width = Pt(1.5)
        shape.line.color.rgb = RGBColor(70, 130, 180)
        shape.shadow.inherit = False

        # Añadir un título moderno dentro de la forma
        title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0), width - Inches(0.6), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        run = title_frame.paragraphs[0].runs[0]
        run.font.size = Pt(self.title_font_size)
        run.font.color.rgb = RGBColor(50, 80, 120)
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(50, 80, 120)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.3),
            top + Inches(1.4),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(70, 130, 180)
        title_underline.line.fill.background()

        # Añadir texto debajo del título con mejor formato
        text_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.5), width - Inches(0.6), height - Inches(1.8))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(50, 50, 70)
        text_frame.margin_top = Inches(-0.25)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.10)
        text_frame.margin_right = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)

        # Elemento decorativo en la esquina inferior derecha
        corner_dec = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + width - Inches(1),
            top + height - Inches(1),
            Inches(0.8),
            Inches(0.8)
        )
        corner_fill = corner_dec.fill
        corner_fill.solid()
        corner_fill.fore_color.rgb = RGBColor(70, 130, 180)
        corner_dec.transparency = 0.6
        corner_dec.line.fill.background()
        corner_dec.rotation = 45

    # Diseño que muestra una imagen oscurecida y un título y contenido a la derecha
    def design4(self, slide, section, content, image_path):
        slide_bg_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_bg_layout)

        # Añadir imagen de fondo completa
        left = Inches(0)
        top = Inches(0)
        width = self.presentation.slide_width
        height = self.presentation.slide_height
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        pic.name = "SLIDE_BACKGROUND_IMAGE"  # Nombre especial para identificarla

        # Añadir una forma de relleno con esquinas redondeadas a la izquierda
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(5)
        height = Inches(5.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 250)
        shape.line.width = Pt(1.5)
        shape.line.color.rgb = RGBColor(0, 130, 114)
        shape.shadow.inherit = False

        # Añadir un título moderno dentro de la forma
        title_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0), width - Inches(0.6), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        run = title_frame.paragraphs[0].runs[0]
        run.font.size = Pt(self.title_font_size)
        run.font.color.rgb = RGBColor(0, 130, 114)
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(0, 130, 114)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(0.3),
            top + Inches(1.2),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(0, 130, 114)
        title_underline.line.fill.background()

        # Añadir texto debajo del título con mejor formato
        text_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.3), width - Inches(0.6), height - Inches(1.6))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(40, 70, 65)
        text_frame.margin_top = Inches(-0.25)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.10)
        text_frame.margin_right = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)

        # Elemento decorativo en la esquina inferior izquierda
        corner_dec = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + width - Inches(1),
            top + height - Inches(1),
            Inches(0.8),
            Inches(0.8)
        )
        corner_fill = corner_dec.fill
        corner_fill.solid()
        corner_fill.fore_color.rgb = RGBColor(0, 130, 114)
        corner_dec.transparency = 0.6
        corner_dec.line.fill.background()
        corner_dec.rotation = 45

    # Diseño que muestra una imagen cuadrada a la izquierda y un título y contenido a la derecha
    def design5(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Añadir fondo con gradiente sutil
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), 
            Inches(0), 
            self.presentation.slide_width, 
            self.presentation.slide_height
        )
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(250, 245, 248)
        bg.line.fill.background()

        # Añadir elemento decorativo - barra lateral
        side_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.3),
            self.presentation.slide_height
        )
        side_bar_fill = side_bar.fill
        side_bar_fill.solid()
        side_bar_fill.fore_color.rgb = RGBColor(175, 0, 80)
        side_bar.line.fill.background()
        
        # Añadir un título con estilo moderno - MODIFICADO para evitar que se sobresalga
        title_left = Inches(0.7)
        title_top = Inches(0.3)
        title_width = self.presentation.slide_width - Inches(1.4)
        title_height = Inches(1)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = title_frame.paragraphs[0].runs[0]
        run.font.size = Pt(self.title_font_size)
        run.font.color.rgb = RGBColor(140, 0, 60)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(140, 0, 60)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.presentation.slide_width/2 - Inches(2),
            Inches(1.6),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(175, 0, 80)
        title_underline.line.fill.background()

        # Añadir una imagen cuadrada a la izquierda - MODIFICADO para usar relleno de forma
        left = Inches(1)
        top = Inches(2)
        side = Inches(4)
        
        # Crear la forma redondeada que contendrá la imagen
        img_container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top,
            side,
            side
        )

        # Configurar el borde de la forma
        img_container.line.width = Pt(2)
        img_container.line.color.rgb = RGBColor(175, 0, 80)
        
        # Añadir la imagen dentro del contenedor
        slide.shapes.add_picture(image_path, left, top, side, side)

        # Añadir un espacio entre la imagen y el texto
        gap = Inches(0.3)
        text_left_bg = img_container.left + img_container.width + gap
        text_top_bg = img_container.top
        text_width_bg = img_container.width
        text_height_bg = img_container.height 

        # Ajustar la posición y tamaño si se sale de la diapositiva
        max_right = self.presentation.slide_width - Inches(0.7)
        if text_left_bg + text_width_bg > max_right:
            text_width_bg = max_right - text_left_bg
            if text_width_bg < Inches(1):
                 text_width_bg = Inches(1)

        text_left_box = text_left_bg + Inches(0.1)
        text_top_box = text_top_bg + Inches(0.1)
        text_width_box = text_width_bg - Inches(0.2)
        text_height_box = text_height_bg - Inches(0.2)
        
        # Fondo para el área de texto
        text_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            text_left_bg,
            text_top_bg,
            text_width_bg,
            text_height_bg
        )
        text_bg_fill = text_bg.fill
        text_bg_fill.solid()
        text_bg_fill.fore_color.rgb = RGBColor(250, 240, 245)
        text_bg.line.width = Pt(1)
        text_bg.line.color.rgb = RGBColor(175, 0, 80)
        text_bg.shadow.inherit = False
        
        # Texto principal
        text_box = slide.shapes.add_textbox(
            text_left_box,
            text_top_box,
            text_width_box,
            text_height_box
        )
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(80, 20, 40)
        text_frame.margin_top = Inches(-0.25)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)
        
        # Elemento decorativo en la esquina inferior derecha
        corner_dec = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.presentation.slide_width - Inches(1.2),
            self.presentation.slide_height - Inches(1.2),
            Inches(1),
            Inches(1)
        )
        corner_fill = corner_dec.fill
        corner_fill.solid()
        corner_fill.fore_color.rgb = RGBColor(175, 0, 80)
        corner_dec.transparency = 0.7
        corner_dec.line.fill.background()
        corner_dec.rotation = 45

    # Diseño que muestra una imagen cuadrada a la derecha y un título y contenido a la izquierda
    def design6(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Añadir fondo con gradiente sutil
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            self.presentation.slide_width,
            self.presentation.slide_height
        )
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(245, 245, 250)
        bg.line.fill.background()

        # Añadir elemento decorativo - barra lateral derecha
        side_bar_left = self.presentation.slide_width - Inches(0.3)
        side_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            side_bar_left,
            Inches(0),
            Inches(0.3),
            self.presentation.slide_height
        )
        side_bar_fill = side_bar.fill
        side_bar_fill.solid()
        side_bar_fill.fore_color.rgb = RGBColor(100, 50, 180)
        side_bar.line.fill.background()

        # Añadir un título con estilo moderno
        title_left = Inches(0.7)
        title_top = Inches(0.3)
        title_width = self.presentation.slide_width - Inches(1.4)
        title_height = Inches(1)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = title_frame.paragraphs[0].runs[0]
        run.font.size = Pt(self.title_font_size)
        run.font.color.rgb = RGBColor(80, 20, 140)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(80, 20, 140)
        
        # Línea decorativa bajo el título
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.presentation.slide_width/2 - Inches(2),
            Inches(1.6),
            Inches(4),
            Inches(0.05)
        )
        underline_fill = title_underline.fill
        underline_fill.solid()
        underline_fill.fore_color.rgb = RGBColor(100, 50, 180)
        title_underline.line.fill.background()

        # Añadir una imagen cuadrada a la derecha
        side = Inches(4)
        right_margin = Inches(0.7)
        img_left = self.presentation.slide_width - side - right_margin
        img_top = Inches(2)

        # Crear la forma rectangular que contendrá la imagen
        img_container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left,
            img_top,
            side,
            side
        )

        # Configurar el borde de la forma
        img_container.line.width = Pt(2)
        img_container.line.color.rgb = RGBColor(100, 50, 180)

        # Añadir la imagen dentro del contenedor
        slide.shapes.add_picture(image_path, img_left, img_top, side, side)

        # Añadir un espacio entre la imagen y el texto
        gap = Inches(0.3)
        
        # Cálculos de posicionamiento espejo para el texto
        text_left_bg = Inches(1)
        text_top_bg = img_top
        text_width_bg = img_left - text_left_bg - gap
        text_height_bg = side

        # Ajustar la posición y tamaño si se sale de la diapositiva
        max_right_text = img_left - gap
        if text_left_bg + text_width_bg > max_right_text:
            text_width_bg = max_right_text - text_left_bg
            if text_width_bg < Inches(1):
                text_width_bg = Inches(1)

        # Márgenes internos
        text_left_box = text_left_bg + Inches(0.1)
        text_top_box = text_top_bg + Inches(0.1)
        text_width_box = text_width_bg - Inches(0.2)
        text_height_box = text_height_bg - Inches(0.2)

        # Fondo para el área de texto
        text_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            text_left_bg,
            text_top_bg,
            text_width_bg,
            text_height_bg
        )
        text_bg_fill = text_bg.fill
        text_bg_fill.solid()
        text_bg_fill.fore_color.rgb = RGBColor(240, 240, 250)
        text_bg.line.width = Pt(1)
        text_bg.line.color.rgb = RGBColor(100, 50, 180)
        text_bg.shadow.inherit = False

        # Texto principal
        text_box = slide.shapes.add_textbox(
            text_left_box,
            text_top_box,
            text_width_box,
            text_height_box
        )
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(60, 20, 120)
        text_frame.margin_top = Inches(-0.25)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)

        # Elemento decorativo en la esquina inferior izquierda
        corner_dec = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.2),
            self.presentation.slide_height - Inches(1.2),
            Inches(1),
            Inches(1)
        )
        corner_fill = corner_dec.fill
        corner_fill.solid()
        corner_fill.fore_color.rgb = RGBColor(100, 50, 180)
        corner_dec.transparency = 0.7
        corner_dec.line.fill.background()
        corner_dec.rotation = 45

    # Diseño que muestra una imagen oscurecida y un título y contenido a la derecha
    def design7(self, slide, section, content, image_path):
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        # Procesar imagen para oscurecerla
        APP_DATA_DIR = os.path.join(os.getenv('APPDATA'), 'Powerpoineador')
        IMAGES_DIR = os.path.join(APP_DATA_DIR, 'images')
        image = Image.open(image_path)
        enhancer = ImageEnhance.Brightness(image)
        image_darker = enhancer.enhance(0.5)
        darker_path = os.path.join(IMAGES_DIR, 'Slide_darker.jpg')
        image_darker.save(darker_path)

        # Añadir imagen oscurecida como fondo
        left = Inches(0)
        top = Inches(0)
        width = self.presentation.slide_width
        height = self.presentation.slide_height
        pic = slide.shapes.add_picture(darker_path, left, top, width, height)
        pic.name = "SLIDE_BACKGROUND_IMAGE"  # Nombre especial para identificarla

        # Añadir título centrado con texto blanco
        title_box = slide.shapes.add_textbox(self.presentation.slide_width / 2 - Inches(2.5), Inches(0.5), Inches(5), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = section
        title_frame.paragraphs[0].runs[0].font.size = Pt(self.title_font_size)
        title_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.margin_bottom = Inches(0)
        title_frame.margin_left = Inches(0.25)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_title_font(title_frame)
        run_title = title_frame.paragraphs[0].runs[0]
        if self.title_bold: run_title.font.bold = True
        if self.title_italic: run_title.font.italic = True
        if self.title_underline: run_title.font.underline = True
        run_title.font.color.rgb = RGBColor(255, 255, 255)

        # Añadir texto centrado con texto blanco
        text_box = slide.shapes.add_textbox(self.presentation.slide_width / 2 - Inches(4), Inches(1.5), Inches(8), Inches(4))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(self.content_font_size)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        text_frame.margin_bottom = Inches(0.08)
        text_frame.margin_left = Inches(0.25)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.apply_content_font(text_frame)