from PySide6.QtWidgets import QWidget, QLabel, QVBoxLayout, QScrollArea, QFrame, QHBoxLayout, QPushButton, QSizePolicy, QApplication, QDialog, QLineEdit, QTextEdit, QFileDialog, QDialogButtonBox, QMessageBox, QProgressDialog, QProgressBar, QComboBox, QCheckBox, QGroupBox, QSpinBox
from PySide6.QtCore import Qt, QSize, QThread, Signal, QObject, QTimer, QRect
from PySide6.QtGui import QPixmap, QFont, QResizeEvent, QIcon, QFontDatabase, QFontMetrics
import os, sys, subprocess
from Traducciones import obtener_traduccion

# Intentar importar python-pptx y manejar el error si no está instalado
try:
    import pptx
    from pptx.dml.color import RGBColor, MSO_COLOR_TYPE
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # Podrías mostrar un mensaje al usuario aquí si es crucial, 
    # pero por ahora solo desactivaremos la funcionalidad.
    print("ADVERTENCIA: La biblioteca 'python-pptx' no está instalada. La edición directa de PPTX estará deshabilitada.")

# Función para obtener la ruta de un recurso
def resource_path(relative_path):
    try:
        # Intenta obtener la ruta base del ejecutable empaquetado
        base_path = sys._MEIPASS
    except Exception:
        # Si no está empaquetado, usa la ruta actual del proyecto
        base_path = os.path.abspath(".")
    
    return os.path.join(base_path, relative_path)

# --- Clase para el diálogo de edición de diapositivas ---
class EditSlideDialog(QDialog):
    def __init__(self, current_data, language, parent=None):
        super().__init__(parent)
        self.current_language = language
        self.parent_widget = parent # VentanaVistaPrevia
        self.temp_preview_updated = False
        
        self.new_image_path = current_data.get('imagen_path') # Inicializar new_image_path con la imagen actual
        self.initial_data = dict(current_data) # Guardar una copia de los datos iniciales

        # Inicializar current_format_settings y original_format_settings aquí
        self.current_format_settings = {} 
        self.original_format_settings = {}

        self.setWindowTitle(obtener_traduccion('edit_slide_dialog_title', self.current_language))
        self.setMinimumWidth(700) # Aumentar ancho mínimo para acomodar más controles

        # --- Layout Principal Horizontal --- 
        main_dialog_layout = QHBoxLayout(self) 
        main_dialog_layout.setContentsMargins(5, 5, 5, 5) # Márgenes pequeños
        main_dialog_layout.setSpacing(10) # Espacio entre botones y contenido

        # --- Botón de Navegación Izquierdo --- 
        self.prev_slide_button = QPushButton()
        left_icon_path = resource_path("iconos/left.png")
        if os.path.exists(left_icon_path):
            self.prev_slide_button.setIcon(QIcon(left_icon_path))
        else:
            self.prev_slide_button.setText("<")
        self.prev_slide_button.setIconSize(QSize(32, 32))
        self.prev_slide_button.setFixedSize(QSize(40, 40))
        self.prev_slide_button.setToolTip(obtener_traduccion('slide_previous', self.current_language))
        self.prev_slide_button.clicked.connect(self.navigate_previous_slide)
        self.prev_slide_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
        self.prev_slide_button.setCursor(Qt.PointingHandCursor)
        main_dialog_layout.addWidget(self.prev_slide_button, 0, Qt.AlignVCenter)

        # --- Contenedor Central Vertical (Contenido del diálogo original) ---
        central_widget = QWidget()
        layout = QVBoxLayout(central_widget) # Layout vertical para el contenido
        layout.setContentsMargins(0,0,0,0)
        layout.setSpacing(5)

        # --- NUEVO: Layout para centrar el botón IA (ARRIBA DEL TODO) ---
        ai_title_button_layout = QHBoxLayout()
        self.generate_ai_title_button = QPushButton()
        ai_title_icon_path = resource_path("iconos/editar_foto.png")
        if os.path.exists(ai_title_icon_path):
            self.generate_ai_title_button.setIcon(QIcon(ai_title_icon_path))
        else:
            self.generate_ai_title_button.setText("IA")
        self.generate_ai_title_button.setIconSize(QSize(24, 24))
        self.generate_ai_title_button.setFixedSize(QSize(32, 32))
        self.generate_ai_title_button.setToolTip(obtener_traduccion('generate_ai_title', self.current_language))
        self.generate_ai_title_button.clicked.connect(self.generate_new_ai_title)
        
        ai_title_button_layout.addStretch(1) # Espaciador izquierdo
        ai_title_button_layout.addWidget(self.generate_ai_title_button)
        ai_title_button_layout.addStretch(1) # Espaciador derecho
        layout.addLayout(ai_title_button_layout) # Añadir layout centrado PRIMERO
        # --- FIN NUEVO Layout ---

        # --- Campos de Edición ---
        # Título (SIN botón IA)
        title_layout = QHBoxLayout()
        title_label = QLabel(obtener_traduccion('edit_title_label', self.current_language))
        self.title_edit = QLineEdit() # Crear vacío, se llenará en load_slide_data
        self.title_edit.textChanged.connect(self.on_text_changed)
        title_layout.addWidget(title_label)
        title_layout.addWidget(self.title_edit)
        layout.addLayout(title_layout) # Añadir DESPUÉS del botón IA

        # --- Contenido: Etiqueta añadida directamente al layout vertical ---
        content_label = QLabel(obtener_traduccion('edit_content_label', self.current_language))
        layout.addWidget(content_label) # Añadir la etiqueta directamente

        self.content_edit = QTextEdit() # Crear vacío, se llenará en load_slide_data
        self.content_edit.setAcceptRichText(False) # Mantener esto para evitar pegado de rich text
        self.content_edit.textChanged.connect(self.on_text_changed)
        self.content_edit.setMinimumHeight(100)
        layout.addWidget(self.content_edit)
        
        # Agregar un pequeño retraso para evitar actualizaciones demasiado frecuentes
        self.text_change_timer = QObject()
        self.text_change_timer.timer = None
        
        # --- NUEVO: Opciones de formato de texto ---
        # Crear dos grupos de opciones: uno para título y otro para contenido
        format_layout = QHBoxLayout()
        
        # --- Grupo de formato para título ---
        title_format_group = QGroupBox(obtener_traduccion('title_font_label', self.current_language))
        title_format_layout = QVBoxLayout(title_format_group)
        
        # Fuente del título
        title_font_layout = QHBoxLayout()
        title_font_layout.addWidget(QLabel(obtener_traduccion('font_label', self.current_language)))
        
        # Usar ComboBox con las mismas fuentes del sistema que PowerpoineatorWidget
        self.title_font_combo = QComboBox()
        # Obtener las fuentes del sistema
        self.system_fonts = QFontDatabase.families()
        if not self.system_fonts:  # Fallback por si no se encuentran fuentes
            self.system_fonts = ["Arial", "Calibri", "Times New Roman"]
            
        # Agregar las fuentes al combobox
        for font_name in self.system_fonts:
            self.title_font_combo.addItem(font_name)
            index = self.title_font_combo.count() - 1
            font = QFont(font_name, 10)
            self.title_font_combo.setItemData(index, font, Qt.FontRole)
            
        # Se configurará con el valor actual en load_format_settings
        self.title_font_combo.currentTextChanged.connect(self.on_title_font_changed)
        title_font_layout.addWidget(self.title_font_combo)
        title_format_layout.addLayout(title_font_layout)
        
        # Tamaño de fuente del título
        title_size_layout = QHBoxLayout()
        title_size_layout.addWidget(QLabel(obtener_traduccion('title_font_size_label', self.current_language)))
        self.title_size_spin = QSpinBox()
        self.title_size_spin.setRange(1, 30)  # Mismo rango que en PowerpoineatorWidget
        # Se configurará con el valor actual en load_format_settings
        self.title_size_spin.valueChanged.connect(self.on_title_size_changed)
        title_size_layout.addWidget(self.title_size_spin)
        title_format_layout.addLayout(title_size_layout)
        
        # Checkboxes para estilo de título
        title_style_layout = QHBoxLayout()
        self.title_bold_check = QCheckBox(obtener_traduccion('title_bold', self.current_language))
        self.title_bold_check.toggled.connect(self.on_title_style_changed)
        title_style_layout.addWidget(self.title_bold_check)
        
        self.title_italic_check = QCheckBox(obtener_traduccion('title_italic', self.current_language))
        self.title_italic_check.toggled.connect(self.on_title_style_changed)
        title_style_layout.addWidget(self.title_italic_check)
        
        self.title_underline_check = QCheckBox(obtener_traduccion('title_underline', self.current_language))
        self.title_underline_check.toggled.connect(self.on_title_style_changed)
        title_style_layout.addWidget(self.title_underline_check)
        
        title_format_layout.addLayout(title_style_layout)
        format_layout.addWidget(title_format_group)
        
        # --- Grupo de formato para contenido ---
        content_format_group = QGroupBox(obtener_traduccion('content_font_label', self.current_language))
        content_format_layout = QVBoxLayout(content_format_group)
        
        # Fuente del contenido
        content_font_layout = QHBoxLayout()
        content_font_layout.addWidget(QLabel(obtener_traduccion('font_label', self.current_language)))
        
        # Usar ComboBox con las mismas fuentes del sistema que PowerpoineatorWidget
        self.content_font_combo = QComboBox()
        # Usar las mismas fuentes obtenidas para el título
        for font_name in self.system_fonts:
            self.content_font_combo.addItem(font_name)
            index = self.content_font_combo.count() - 1
            font = QFont(font_name, 10)
            self.content_font_combo.setItemData(index, font, Qt.FontRole)
            
        # Se configurará con el valor actual en load_format_settings
        self.content_font_combo.currentTextChanged.connect(self.on_content_font_changed)
        content_font_layout.addWidget(self.content_font_combo)
        content_format_layout.addLayout(content_font_layout)
        
        # Tamaño de fuente del contenido
        content_size_layout = QHBoxLayout()
        content_size_layout.addWidget(QLabel(obtener_traduccion('content_font_size_label', self.current_language)))
        self.content_size_spin = QSpinBox()
        self.content_size_spin.setRange(1, 20)  # Mismo rango que en PowerpoineatorWidget
        # Se configurará con el valor actual en load_format_settings
        self.content_size_spin.valueChanged.connect(self.on_content_size_changed)
        content_size_layout.addWidget(self.content_size_spin)
        content_format_layout.addLayout(content_size_layout)
        
        # Checkboxes para estilo de contenido
        content_style_layout = QHBoxLayout()
        self.content_bold_check = QCheckBox(obtener_traduccion('content_bold', self.current_language))
        self.content_bold_check.toggled.connect(self.on_content_style_changed)
        content_style_layout.addWidget(self.content_bold_check)
        
        self.content_italic_check = QCheckBox(obtener_traduccion('content_italic', self.current_language))
        self.content_italic_check.toggled.connect(self.on_content_style_changed)
        content_style_layout.addWidget(self.content_italic_check)
        
        self.content_underline_check = QCheckBox(obtener_traduccion('content_underline', self.current_language))
        self.content_underline_check.toggled.connect(self.on_content_style_changed)
        content_style_layout.addWidget(self.content_underline_check)
        
        content_format_layout.addLayout(content_style_layout)
        format_layout.addWidget(content_format_group)
        
        # Añadir el layout de formato al layout principal
        layout.addLayout(format_layout)
        # --- FIN NUEVO: Opciones de formato de texto ---
        
        # Imagen
        image_layout = QHBoxLayout()
        image_label_text = obtener_traduccion('edit_image_label', self.current_language)
        image_label = QLabel(image_label_text)
        self.current_image_label = QLabel("") # Crear vacío, se llenará en load_slide_data
        self.current_image_label.setWordWrap(True)
        browse_button = QPushButton(obtener_traduccion('edit_browse_button', self.current_language))
        browse_button.clicked.connect(self.browse_new_image)
        
        self.generate_ai_button = QPushButton(QIcon(resource_path("iconos/editar_foto.png")), "") 
        self.generate_ai_button.setIconSize(QSize(24, 24))
        self.generate_ai_button.setFixedSize(QSize(32, 32))
        self.generate_ai_button.setToolTip(obtener_traduccion('generate_ai_image', self.current_language))
        self.generate_ai_button.clicked.connect(self.generate_new_ai_image)
        
        image_layout.addWidget(image_label)
        image_layout.addWidget(self.current_image_label, 1) # Darle más espacio
        image_layout.addWidget(browse_button)
        image_layout.addWidget(self.generate_ai_button)
        layout.addLayout(image_layout)

        # --- Botones OK/Cancel --- 
        button_box = QDialogButtonBox(QDialogButtonBox.Save | QDialogButtonBox.Cancel)
        save_button = button_box.button(QDialogButtonBox.Save)
        if save_button:
            save_button.setText(obtener_traduccion('edit_save', self.current_language))
        cancel_button = button_box.button(QDialogButtonBox.Cancel)
        if cancel_button:
            cancel_button.setText(obtener_traduccion('edit_cancel', self.current_language))
            
        button_box.accepted.connect(self.accept) # Guardar cambios de la diapositiva ACTUAL
        button_box.rejected.connect(self.reject)
        layout.addWidget(button_box)
        
        # --- Añadir widget central al layout principal --- 
        main_dialog_layout.addWidget(central_widget, 1) # Darle más espacio horizontal que a los botones

        # --- Botón de Navegación Derecho --- 
        self.next_slide_button = QPushButton()
        right_icon_path = resource_path("iconos/right.png")
        if os.path.exists(right_icon_path):
            self.next_slide_button.setIcon(QIcon(right_icon_path))
        else:
            self.next_slide_button.setText(">")
        self.next_slide_button.setIconSize(QSize(32, 32))
        self.next_slide_button.setFixedSize(QSize(40, 40))
        self.next_slide_button.setToolTip(obtener_traduccion('slide_next', self.current_language))
        self.next_slide_button.clicked.connect(self.navigate_next_slide)
        self.next_slide_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
        self.next_slide_button.setCursor(Qt.PointingHandCursor)
        main_dialog_layout.addWidget(self.next_slide_button, 0, Qt.AlignVCenter)

        # --- Inicialización Final --- 
        self.load_format_settings() # Cargar configuración de formato primero
        self.load_slide_data() # Cargar datos de la diapositiva inicial
        self.update_navigation_button_state() # Establecer estado inicial de botones
        self.finished.connect(self.on_dialog_finished) # Conexión para restaurar vista previa al cerrar

    # >>> NUEVO: Método para cargar la configuración de formato <<<
    def load_format_settings(self):
        # Prioridad 1: Usar la configuración específica de la diapositiva si existe y fue guardada
        slide_specific_format = self.initial_data.get('format_settings')

        if slide_specific_format:
            # Si la diapositiva tiene su propio formato guardado, usarlo como base para el diálogo
            self.original_format_settings = dict(slide_specific_format)
            print(f"Diálogo: Cargando formato específico de la diapositiva: {self.original_format_settings}")
        elif self.parent_widget and hasattr(self.parent_widget, 'title_font_name'):
            # Prioridad 2: Si no hay formato específico, usar el formato global de VentanaVistaPrevia
            self.original_format_settings = {
                'title_font_name': self.parent_widget.title_font_name,
                'content_font_name': self.parent_widget.content_font_name,
                'title_font_size': self.parent_widget.title_font_size,
                'content_font_size': self.parent_widget.content_font_size,
                'title_bold': self.parent_widget.title_bold,
                'title_italic': self.parent_widget.title_italic,
                'title_underline': self.parent_widget.title_underline,
                'content_bold': self.parent_widget.content_bold,
                'content_italic': self.parent_widget.content_italic,
                'content_underline': self.parent_widget.content_underline,
            }
            print(f"Diálogo: Cargando formato global de VentanaVistaPrevia: {self.original_format_settings}")
        else:
            # Prioridad 3: Fallback a valores predeterminados si no hay nada más
            print("Advertencia: Faltan atributos de formato en parent_widget. Usando predeterminados para el diálogo.")
            default_font = "Calibri"
            default_title_size = 24
            default_content_size = 12
            self.original_format_settings = {
                'title_font_name': default_font, 'content_font_name': default_font,
                'title_font_size': default_title_size, 'content_font_size': default_content_size,
                'title_bold': False, 'title_italic': False, 'title_underline': False,
                'content_bold': False, 'content_italic': False, 'content_underline': False,
            }

        # current_format_settings siempre empieza como una copia de original_format_settings al cargar el diálogo o la diapositiva
        self.current_format_settings = dict(self.original_format_settings)
        
        # Configurar controles con los valores actuales (de current_format_settings)
        # (Bloquear señales temporalmente para evitar llamadas a update_preview_temporarily)
        self.title_font_combo.blockSignals(True)
        self.title_size_spin.blockSignals(True)
        self.title_bold_check.blockSignals(True)
        self.title_italic_check.blockSignals(True)
        self.title_underline_check.blockSignals(True)
        self.content_font_combo.blockSignals(True)
        self.content_size_spin.blockSignals(True)
        self.content_bold_check.blockSignals(True)
        self.content_italic_check.blockSignals(True)
        self.content_underline_check.blockSignals(True)

        title_font_index = self.title_font_combo.findText(self.current_format_settings['title_font_name'])
        self.title_font_combo.setCurrentIndex(title_font_index if title_font_index >=0 else self.title_font_combo.findText("Calibri"))
        
        self.title_size_spin.setValue(self.current_format_settings['title_font_size'])
        self.title_bold_check.setChecked(self.current_format_settings['title_bold'])
        self.title_italic_check.setChecked(self.current_format_settings['title_italic'])
        self.title_underline_check.setChecked(self.current_format_settings['title_underline'])
        
        content_font_index = self.content_font_combo.findText(self.current_format_settings['content_font_name'])
        self.content_font_combo.setCurrentIndex(content_font_index if content_font_index >=0 else self.content_font_combo.findText("Calibri"))

        self.content_size_spin.setValue(self.current_format_settings['content_font_size'])
        self.content_bold_check.setChecked(self.current_format_settings['content_bold'])
        self.content_italic_check.setChecked(self.current_format_settings['content_italic'])
        self.content_underline_check.setChecked(self.current_format_settings['content_underline'])
        
        # Reactivar señales
        self.title_font_combo.blockSignals(False)
        self.title_size_spin.blockSignals(False)
        self.title_bold_check.blockSignals(False)
        self.title_italic_check.blockSignals(False)
        self.title_underline_check.blockSignals(False)
        self.content_font_combo.blockSignals(False)
        self.content_size_spin.blockSignals(False)
        self.content_bold_check.blockSignals(False)
        self.content_italic_check.blockSignals(False)
        self.content_underline_check.blockSignals(False)

    # >>> NUEVO: Métodos para manejar cambios en los controles de formato <<<
    def on_title_font_changed(self, font):
        self.current_format_settings['title_font_name'] = font
        self.update_preview_temporarily()
        
    def on_title_size_changed(self, value):
        self.current_format_settings['title_font_size'] = value
        self.update_preview_temporarily()
        
    def on_title_style_changed(self):
        self.current_format_settings['title_bold'] = self.title_bold_check.isChecked()
        self.current_format_settings['title_italic'] = self.title_italic_check.isChecked()
        self.current_format_settings['title_underline'] = self.title_underline_check.isChecked()
        self.update_preview_temporarily()
        
    def on_content_font_changed(self, font):
        self.current_format_settings['content_font_name'] = font
        self.update_preview_temporarily()
        
    def on_content_size_changed(self, value):
        self.current_format_settings['content_font_size'] = value
        self.update_preview_temporarily()
        
    def on_content_style_changed(self):
        self.current_format_settings['content_bold'] = self.content_bold_check.isChecked()
        self.current_format_settings['content_italic'] = self.content_italic_check.isChecked()
        self.current_format_settings['content_underline'] = self.content_underline_check.isChecked()
        self.update_preview_temporarily()

    # >>> NUEVO: Cargar datos de la diapositiva actual en los campos <<< 
    def load_slide_data(self):
        if not self.parent_widget or not hasattr(self.parent_widget, 'all_slides_data') or not self.parent_widget.all_slides_data:
            return
        
        current_index = self.parent_widget.current_slide_index
        if 0 <= current_index < len(self.parent_widget.all_slides_data):
            slide_data = self.parent_widget.all_slides_data[current_index]
            
            self.initial_data = dict(slide_data) 
            self.new_image_path = self.initial_data.get('imagen_path') # Asegurar que new_image_path se actualiza
            
            self.title_edit.setText(self.initial_data.get('titulo', ''))
            self.content_edit.setPlainText(self.initial_data.get('contenido', ''))
            self.current_image_label.setText(os.path.basename(self.new_image_path) if self.new_image_path else "N/A")
            
            # Cargar la configuración de formato específica de esta diapositiva (o la global si no tiene)
            self.load_format_settings() # <--- Llamar aquí para cargar el formato de la diapositiva actual

            self.temp_preview_updated = False 
        else:
            self.title_edit.clear()
            self.content_edit.clear()
            self.new_image_path = None
            self.current_image_label.setText("N/A")
            self.initial_data = {}
            # Incluso si no hay datos, cargar un formato predeterminado
            self.load_format_settings() # <--- Llamar aquí también para estado vacío

    # >>> NUEVO: Actualizar estado de los botones de navegación <<< 
    def update_navigation_button_state(self):
        if not self.parent_widget or not hasattr(self.parent_widget, 'all_slides_data') or not self.parent_widget.all_slides_data:
            self.prev_slide_button.setEnabled(False)
            self.next_slide_button.setEnabled(False)
            return
        
        current_index = self.parent_widget.current_slide_index
        total_slides = len(self.parent_widget.all_slides_data)
        
        self.prev_slide_button.setEnabled(current_index > 0)
        self.next_slide_button.setEnabled(current_index < total_slides - 1)

    # >>> NUEVO: Navegar a la diapositiva anterior <<< 
    def navigate_previous_slide(self):
        if not self.parent_widget:
            return
        current_index = self.parent_widget.current_slide_index
        if current_index > 0:
            if self.temp_preview_updated:
                 # Restaurar datos de la diapositiva que se deja
                 self.parent_widget.all_slides_data[current_index] = dict(self.initial_data)
                 # Restaurar el formato global de VentanaVistaPrevia al que tenía ANTES de editar ESTA diapositiva
                 self.parent_widget.update_format_settings(**self.original_format_settings) 
                 self.parent_widget.mostrar_diapositiva_actual() # Mostrar brevemente restaurada
                 QApplication.processEvents()

            self.parent_widget.mostrar_diapositiva_anterior()
            self.load_slide_data() # Esto cargará los datos Y el formato de la nueva diapositiva
            self.update_navigation_button_state()

    # >>> NUEVO: Navegar a la diapositiva siguiente <<< 
    def navigate_next_slide(self):
        if not self.parent_widget or not hasattr(self.parent_widget, 'all_slides_data'):
            return
            
        current_index = self.parent_widget.current_slide_index
        total_slides = len(self.parent_widget.all_slides_data)
        if current_index < total_slides - 1:
            if self.temp_preview_updated:
                 self.parent_widget.all_slides_data[current_index] = dict(self.initial_data)
                 self.parent_widget.update_format_settings(**self.original_format_settings)
                 self.parent_widget.mostrar_diapositiva_actual()
                 QApplication.processEvents()

            self.parent_widget.mostrar_diapositiva_siguiente() 
            self.load_slide_data() # Esto cargará los datos Y el formato de la nueva diapositiva
            self.update_navigation_button_state()

    # >>> INICIO: Funciones para generar título con IA (SIN CAMBIOS EN ESTA PARTE) <<<
    def _find_main_powerpoineator_widget(self):
        """Función auxiliar para encontrar el widget principal PowerpoineatorWidget."""
        main_widget = None
        if not self.parent_widget:
            return None

        # Estrategia 1: El padre directo es el widget principal
        if hasattr(self.parent_widget, 'descripcion_text') and \
           hasattr(self.parent_widget, 'texto_combo') and \
           hasattr(self.parent_widget, 'num_diapositivas_spin') and \
           hasattr(self.parent_widget, 'current_language'):
            main_widget = self.parent_widget
            return main_widget

        # Estrategia 2: El "abuelo" es el widget principal
        parent_obj = self.parent_widget.parent()
        if parent_obj and hasattr(parent_obj, 'descripcion_text') and \
           hasattr(parent_obj, 'texto_combo') and \
           hasattr(parent_obj, 'num_diapositivas_spin') and \
           hasattr(parent_obj, 'current_language'):
            main_widget = parent_obj
            return main_widget
        
        # Estrategia 3: Ir un nivel más arriba si el abuelo no lo era pero existe un bisabuelo
        if parent_obj and hasattr(parent_obj, 'parent'):
            grandparent_obj = parent_obj.parent()
            if grandparent_obj and hasattr(grandparent_obj, 'descripcion_text') and \
               hasattr(grandparent_obj, 'texto_combo') and \
               hasattr(grandparent_obj, 'num_diapositivas_spin') and \
               hasattr(grandparent_obj, 'current_language'):
                main_widget = grandparent_obj
                return main_widget

        # Estrategia 4: Buscar en todos los widgets de nivel superior
        for widget_app in QApplication.topLevelWidgets():
            if hasattr(widget_app, 'widget') and widget_app.widget and \
               hasattr(widget_app.widget, 'descripcion_text') and \
               hasattr(widget_app.widget, 'texto_combo') and \
               hasattr(widget_app.widget, 'num_diapositivas_spin') and \
               hasattr(widget_app.widget, 'current_language'):
                main_widget = widget_app.widget
                return main_widget
        
        return None # No encontrado o no válido

    def generate_new_ai_title(self):
        from PySide6.QtWidgets import QApplication, QStyle
        from PySide6.QtCore import Qt
        
        self.title_progress_dialog = QProgressDialog(
            obtener_traduccion('generate_ai_title_progress_label', self.current_language),
            obtener_traduccion('edit_cancel', self.current_language), 0, 0, self
        )
        self.title_progress_dialog.setWindowTitle(obtener_traduccion('generate_ai_title_progress_title', self.current_language))
        self.title_progress_dialog.setWindowModality(Qt.WindowModal)
        self.title_progress_dialog.setMinimumDuration(0)
        self.title_progress_dialog.setMinimumWidth(450)
        self.title_progress_dialog.setAutoClose(False)
        self.title_progress_dialog.setAutoReset(False)

        title_progress_bar_widget = self.title_progress_dialog.findChild(QProgressBar)
        if title_progress_bar_widget:
            title_progress_bar_widget.setMinimumWidth(400)
            title_progress_bar_widget.setMaximumWidth(420)
            title_progress_bar_widget.setAlignment(Qt.AlignCenter)
        else:
            print("Advertencia: No se pudo encontrar QProgressBar en QProgressDialog para el título.")

        desktop = QApplication.primaryScreen().availableGeometry()
        dialog_rect = QStyle.alignedRect(Qt.LeftToRight, Qt.AlignCenter, self.title_progress_dialog.size(), desktop)
        self.title_progress_dialog.setGeometry(dialog_rect)
        
        self.title_progress_dialog.show()
        QApplication.processEvents() # Clave para que la UI se actualice y la barra de progreso empiece a animar
        
        slide_idx_for_this_request = self.parent_widget.current_slide_index if hasattr(self.parent_widget, 'current_slide_index') else 0

        # --- Recopilación RÁPIDA de referencias en el hilo principal ---
        main_powerpoineator_widget = self._find_main_powerpoineator_widget()

        if not main_powerpoineator_widget:
            self.title_progress_dialog.close()
            # Usar traducciones existentes o genéricas si las nuevas no están disponibles aún
            error_title = obtener_traduccion('error', self.current_language)
            error_msg = obtener_traduccion('ai_title_main_window_access_error', self.current_language) # Buscar una traducción existente adecuada
            if error_msg == 'ai_title_main_window_access_error': error_msg = "No se puede acceder a la ventana principal para obtener datos."
            QMessageBox.warning(self, error_title, error_msg)
            return

        required_attrs = ['descripcion_text', 'texto_combo', 'num_diapositivas_spin', 'current_language']
        missing_attrs = [attr for attr in required_attrs if not hasattr(main_powerpoineator_widget, attr)]
        if missing_attrs:
            self.title_progress_dialog.close()
            error_title = obtener_traduccion('error', self.current_language)
            error_msg_format = obtener_traduccion('ai_title_missing_components_error', self.current_language) # Buscar una traducción existente adecuada
            if error_msg_format == 'ai_title_missing_components_error': error_msg_format = "Faltan componentes necesarios en la ventana principal: {components}"
            QMessageBox.warning(self, error_title, error_msg_format.format(components=', '.join(missing_attrs)))
            return
        
        class GenerateTitleSignals(QObject):
            success = Signal(list)
            error = Signal(str)
            finished = Signal()

        class GenerateTitleWorker(QThread):
            def __init__(self, main_widget_ref, dialog_lang_ref, slide_idx_ref):
                super().__init__()
                self.main_widget = main_widget_ref
                self.dialog_language = dialog_lang_ref
                self.slide_idx = slide_idx_ref
                self.signals = GenerateTitleSignals()
                self._is_interrupted = False

            def run(self):
                try:
                    descripcion_original = self.main_widget.descripcion_text.toPlainText()
                    modelo_texto = self.main_widget.texto_combo.currentText()
                    num_diapositivas_original = self.main_widget.num_diapositivas_spin.value()
                    main_app_language = self.main_widget.current_language

                    if not modelo_texto:
                        error_msg = obtener_traduccion('ai_title_missing_model_message', self.dialog_language)
                        self.signals.error.emit(error_msg)
                        return
                    
                    instruccion_idioma = obtener_traduccion('language_instruction', main_app_language)
                    full_description_for_ia = descripcion_original
                    
                    pdf_file_path = getattr(self.main_widget, 'pdf_cargado', None)
                    if pdf_file_path and os.path.exists(pdf_file_path):
                        try:
                            texto_doc = ""
                            if hasattr(self.main_widget, 'extraer_texto_pdf') and hasattr(self.main_widget, 'extraer_texto_txt'):
                                if pdf_file_path.lower().endswith('.pdf'):
                                    texto_doc = self.main_widget.extraer_texto_pdf(pdf_file_path)
                                elif pdf_file_path.lower().endswith('.txt'):
                                    texto_doc = self.main_widget.extraer_texto_txt(pdf_file_path)
                            
                            if texto_doc.strip():
                                # Escapar correctamente las nuevas líneas para f-string
                                full_description_for_ia += f"\n\nContenido del documento para tener en cuenta:\n{texto_doc}"
                        except Exception as e_doc:
                            print(f"Error al procesar el documento para el título AI en worker: {str(e_doc)}")

                    final_prompt_descripcion = full_description_for_ia + f" hazlo OBLIGATORIAMENTE en {num_diapositivas_original} claves-valores (diapositivas) y en {instruccion_idioma}, tal que los títulos de la tupla NO superen 4 palabras y del contenido NO superen 69 palabras"

                    if self.isInterruptionRequested(): return

                    from Logica_diapositivas import obtener_respuesta_ia # Importar aquí para evitar problemas de importación cíclica a nivel de módulo
                    tupla_respuesta_dict = obtener_respuesta_ia(final_prompt_descripcion, modelo_texto, None)

                    if self.isInterruptionRequested(): return

                    if tupla_respuesta_dict and isinstance(tupla_respuesta_dict, dict):
                        lista_ordenada_items = list(tupla_respuesta_dict.items())
                        self.signals.success.emit(lista_ordenada_items)
                    else:
                        error_msg = obtener_traduccion('ai_title_invalid_response', self.dialog_language) # Reutilizar traducción existente
                        if error_msg == 'ai_title_invalid_response': error_msg = "La respuesta del modelo no tiene un formato válido o está vacía."
                        self.signals.error.emit(error_msg)
                except ImportError as e_imp:
                    error_msg = obtener_traduccion('ai_module_import_error', self.dialog_language) # Buscar traducción adecuada
                    if error_msg == 'ai_module_import_error': error_msg = f"Error de importación en el worker: {e_imp}"
                    self.signals.error.emit(error_msg)
                except Exception as e:
                    self.signals.error.emit(str(e))
                finally:
                    self.signals.finished.emit()
            
            def requestInterruption(self):
                self._is_interrupted = True
                super().requestInterruption()
            
            def isInterruptionRequested(self):
                return self._is_interrupted or super().isInterruptionRequested()

        self.title_worker = GenerateTitleWorker(main_powerpoineator_widget, 
                                                self.current_language,
                                                slide_idx_for_this_request)
                                                
        self.title_worker.signals.success.connect(
            lambda lista_items: self.handle_title_success(lista_items, slide_idx_for_this_request)
        )
        self.title_worker.signals.error.connect(self.handle_title_error) 
        self.title_worker.signals.finished.connect(self.handle_title_finished)
        self.title_progress_dialog.canceled.connect(self.cancel_title_generation)
        
        self.title_worker.start()

    def handle_title_success(self, lista_ordenada_items, captured_slide_idx):
        if hasattr(self, 'title_progress_dialog') and self.title_progress_dialog.isVisible():
            self.title_progress_dialog.close()
        try:
            current_slide_idx = captured_slide_idx

            if current_slide_idx < 0:
                QMessageBox.warning(self, obtener_traduccion('error', self.current_language), 
                                    obtener_traduccion('ai_title_invalid_slide_index', self.current_language))
                return
            
            if not isinstance(lista_ordenada_items, list) or not lista_ordenada_items:
                 QMessageBox.warning(self, obtener_traduccion('error', self.current_language),
                                     obtener_traduccion('ai_title_invalid_response', self.current_language))
                 return

            if current_slide_idx < len(lista_ordenada_items):
                nuevo_titulo = lista_ordenada_items[current_slide_idx][0]
                nuevo_contenido = lista_ordenada_items[current_slide_idx][1]
                
                self.title_edit.setText(nuevo_titulo)
                self.content_edit.setPlainText(nuevo_contenido)
                
                print(f"INFO: Título actualizado a: {nuevo_titulo}")
                print(f"INFO: Contenido actualizado a: {nuevo_contenido}")

                QMessageBox.information(self,
                                        obtener_traduccion('ai_title_success_title', self.current_language),
                                        obtener_traduccion('ai_title_success_message', self.current_language))
            else:
                QMessageBox.warning(self,
                                    obtener_traduccion('ai_title_index_error_title', self.current_language),
                                    obtener_traduccion('ai_title_index_error_message', self.current_language).format(index=current_slide_idx + 1, total=len(lista_ordenada_items)))
        except Exception as e:
            QMessageBox.critical(self, obtener_traduccion('error', self.current_language),
                                 f"{obtener_traduccion('ai_title_processing_error', self.current_language)}: {str(e)}")

    def handle_title_error(self, error_msg):
        if hasattr(self, 'title_progress_dialog') and self.title_progress_dialog and self.title_progress_dialog.isVisible():
            self.title_progress_dialog.close()
        QMessageBox.critical(self,
                             obtener_traduccion('ai_title_error_title', self.current_language),
                             obtener_traduccion('ai_title_error_message', self.current_language).format(error=error_msg))

    def handle_title_finished(self):
        # Asegurarse de que el diálogo se cierre si el worker termina y el diálogo sigue visible
        # (aunque success y error ya deberían haberlo cerrado)
        if hasattr(self, 'title_progress_dialog') and self.title_progress_dialog and self.title_progress_dialog.isVisible():
             self.title_progress_dialog.close()
        if hasattr(self, 'title_worker') and self.title_worker:
            self.title_worker.deleteLater()
            try:
                del self.title_worker 
            except AttributeError:
                pass # Puede que ya se haya eliminado si hubo un error rápido o cancelación

    def cancel_title_generation(self):
        if hasattr(self, 'title_worker') and self.title_worker and self.title_worker.isRunning():
            self.title_worker.requestInterruption()
    # >>> FIN: Nuevas funciones para generar título con IA <<<

    def on_text_changed(self):
        """Manejar cambios en el texto con un pequeño retraso para no actualizar constantemente"""
        # Cancelar timer anterior si existe
        if hasattr(self.text_change_timer, 'timer') and self.text_change_timer.timer:
            self.text_change_timer.timer.stop()
        
        # Crear un nuevo timer
        self.text_change_timer.timer = QTimer()
        self.text_change_timer.timer.setSingleShot(True)
        self.text_change_timer.timer.timeout.connect(self.update_preview_temporarily)
        self.text_change_timer.timer.start(0)

    def browse_new_image(self):
        image_filter_key = 'edit_image_filter'
        image_filter_formats = "*.png *.jpg *.jpeg *.bmp *.gif *.tiff"
        image_filter_text = obtener_traduccion(image_filter_key, self.current_language)
        if image_filter_text == image_filter_key: # Fallback si la traducción no existe
            image_filter_text = f"Imágenes ({image_filter_formats})"
        else:
            try:
                image_filter_text = image_filter_text.format(image_filter_formats)
            except KeyError:
                # Si el formato no funciona (p.ej., {} no presente en la traducción)
                image_filter_text = f"Imágenes ({image_filter_formats})"
                
        select_image_title = obtener_traduccion('edit_select_image', self.current_language)

        file_path, _ = QFileDialog.getOpenFileName(
            self,
            select_image_title,
            "", # Directorio inicial
            image_filter_text
        )
        if file_path and os.path.exists(file_path):
            self.new_image_path = file_path
            self.current_image_label.setText(os.path.basename(file_path))
            # Actualizar vista previa en tiempo real
            self.update_preview_temporarily()
            
    def generate_new_ai_image(self):
        """Generar una nueva imagen usando el modelo de IA seleccionado en el combobox"""
        if not self.parent_widget or not hasattr(self.parent_widget, 'parent'):
            QMessageBox.warning(self, "Error", "No se puede acceder a la ventana principal")
            return
            
        # Buscar el widget principal para obtener el modelo de imagen seleccionado
        main_widget = self.parent_widget.parent()
        while main_widget and not hasattr(main_widget, 'imagen_combo'):
            if hasattr(main_widget, 'parent'):
                main_widget = main_widget.parent()
            else:
                main_widget = None
                
        if not main_widget:
            QMessageBox.warning(self,
                                obtener_traduccion('ai_image_access_error_title', self.current_language),
                                obtener_traduccion('ai_image_access_main_window_message', self.current_language))
            return
        if not hasattr(main_widget, 'imagen_combo'):
            QMessageBox.warning(self,
                                obtener_traduccion('ai_image_access_error_title', self.current_language),
                                obtener_traduccion('ai_image_access_model_message', self.current_language))
            return
            
        # Obtener modelo de imagen y otros datos necesarios
        modelo_imagen = main_widget.imagen_combo.currentText()
        titulo = self.title_edit.text()
        contenido = self.content_edit.toPlainText()
        
        if not titulo or not contenido:
            QMessageBox.warning(self,
                                obtener_traduccion('ai_image_missing_data_title', self.current_language),
                                obtener_traduccion('ai_image_missing_data_message', self.current_language))
            return
            
        # Verificar si necesitamos imagen personalizada para el modelo
        imagen_personalizada = None
        if modelo_imagen in ['flux-pulid [$0.029]', 'photomaker [$0.0069]']:
            if hasattr(main_widget, 'imagen_personalizada') and main_widget.imagen_personalizada:
                imagen_personalizada = main_widget.imagen_personalizada
            else:
                QMessageBox.warning(self,
                                    obtener_traduccion('ai_image_missing_custom_image_title', self.current_language),
                                    obtener_traduccion('ai_image_missing_custom_image_message', self.current_language))
                return
                
        # Crear un diálogo de progreso
        class GenerateImageSignals(QObject):
            success = Signal(str)
            error = Signal(str)
            finished = Signal()
        
        class GenerateImageWorker(QThread):
            def __init__(self, modelo_imagen, titulo, contenido, imagen_personalizada):
                super().__init__()
                self.modelo_imagen = modelo_imagen
                self.titulo = titulo
                self.contenido = contenido
                self.imagen_personalizada = imagen_personalizada
                self.signals = GenerateImageSignals()
                self.temp_imagen_path = None
                # --- NUEVO: Flag para interrupción ---
                self._is_interrupted = False
                
            def run(self):
                try:
                    # Importar funciones necesarias
                    from Logica_diapositivas import generar_imagen_ia, generar_imagen_flux, generar_imagen_photomaker
                    import sys, os
                    
                    # --- NUEVO: Comprobar interrupción temprana ---
                    if self.isInterruptionRequested():
                        print("Generación de imagen cancelada antes de empezar.")
                        return
                        
                    # Definir directorios para imágenes temporales
                    if sys.platform == 'win32':
                        APP_DATA_DIR = os.path.join(os.getenv('APPDATA'), 'Powerpoineador')
                    elif sys.platform == 'darwin':
                        APP_DATA_DIR = os.path.join(os.path.expanduser('~'), 'Library', 'Application Support', 'Powerpoineador')
                    else:
                        APP_DATA_DIR = os.path.join(os.path.expanduser('~'), '.Powerpoineador')
                        
                    IMAGES_DIR = os.path.join(APP_DATA_DIR, 'images')
                    os.makedirs(IMAGES_DIR, exist_ok=True)
                    
                    # Generar una imagen temporal con nombre único basado en timestamp
                    import time
                    timestamp = int(time.time())
                    self.temp_imagen_path = os.path.join(IMAGES_DIR, f'temp_edit_slide_{timestamp}.jpg')
                    
                    # --- NUEVO: Comprobar interrupción antes de la llamada costosa ---
                    if self.isInterruptionRequested():
                        print("Generación de imagen cancelada antes de la llamada a la IA.")
                        return
                        
                    # Generar la imagen según el modelo
                    if self.modelo_imagen in ['flux-pulid [$0.029]', 'photomaker [$0.0069]']:
                        if self.modelo_imagen == 'flux-pulid [$0.029]':
                            img = generar_imagen_flux(self.titulo, self.contenido, "", self.imagen_personalizada, None, False)
                        else:  # photomaker
                            img = generar_imagen_photomaker(self.titulo, self.contenido, "", self.imagen_personalizada, None, False)
                    else:
                        img = generar_imagen_ia(self.titulo, self.contenido, "", self.modelo_imagen, None)
                        
                    # Guardar la imagen generada
                    img.save(self.temp_imagen_path)
                    
                    # --- NUEVO: Comprobar interrupción antes de emitir señal ---
                    if self.isInterruptionRequested():
                        print("Generación de imagen cancelada después de generar pero antes de emitir señal.")
                        # Opcional: Eliminar imagen temporal si se cancela aquí
                        if os.path.exists(self.temp_imagen_path):
                             try:
                                 os.remove(self.temp_imagen_path)
                                 print(f"Imagen temporal eliminada: {self.temp_imagen_path}")
                             except OSError as e:
                                 print(f"Error eliminando imagen temporal: {e}")
                        return
                        
                    # Emitir señal de éxito con la ruta de la imagen
                    self.signals.success.emit(self.temp_imagen_path)
                    
                except Exception as e:
                    # Emitir señal de error
                    self.signals.error.emit(str(e))
                finally:
                    # Emitir señal de finalización
                    self.signals.finished.emit()
        
        # Crear y configurar el diálogo de progreso - MOSTRAR PRIMERO
        progress_dialog = QProgressDialog(
            obtener_traduccion('generate_ai_image_label', self.current_language), # Usar traducción para etiqueta
            obtener_traduccion('edit_cancel', self.current_language), # Usar traducción para botón cancelar
            0, 0, self
        )
        progress_dialog.setWindowTitle(obtener_traduccion('generate_ai_image_title', self.current_language)) # Usar traducción para título
        progress_dialog.setWindowModality(Qt.WindowModal)
        progress_dialog.setMinimumDuration(0) # Mostrar inmediatamente
        progress_dialog.setMinimumWidth(450) # Aumentar ancho mínimo
        progress_dialog.setAutoClose(False) # Evitar que se cierre automáticamente
        progress_dialog.setAutoReset(False) # Evitar que se reinicie automáticamente
        
        # --- MODIFICADO: Encontrar QProgressBar y aplicar estilo ---
        # Buscar el QProgressBar dentro del diálogo
        progress_bar_widget = progress_dialog.findChild(QProgressBar) # Buscar QProgressBar directamente

        if progress_bar_widget:
            # Hacer que la barra de progreso ocupe más espacio horizontalmente
            # Establecer un tamaño mínimo más grande relativo al diálogo
            progress_bar_widget.setMinimumWidth(400)
            progress_bar_widget.setMaximumWidth(420) # Limitar para que no sea *exactamente* el ancho del diálogo
            progress_bar_widget.setAlignment(Qt.AlignCenter) # Centrar texto si lo hubiera (aunque en modo indeterminado no suele haber)

            # --- MODIFICADO: Eliminar hoja de estilos para usar el estilo del sistema ---
            # Al no establecer una hoja de estilos, QProgressBar usará la apariencia nativa.
            # progress_bar_widget.setStyleSheet(""" ... """) # <- Eliminado
            # --- FIN MODIFICACIÓN ESTILO ---

        else:
            print("Advertencia: No se pudo encontrar QProgressBar en QProgressDialog para aplicar estilo.")
        # --- FIN MODIFICACIÓN ---

        # NUEVO: Centrar el diálogo en la pantalla usando QStyle
        from PySide6.QtWidgets import QStyle
        from PySide6.QtCore import QRect
        desktop = QApplication.primaryScreen().availableGeometry()
        dialog_rect = QStyle.alignedRect(
            Qt.LeftToRight,
            Qt.AlignCenter,
            progress_dialog.size(),
            desktop
        )
        progress_dialog.setGeometry(dialog_rect)
        
        # Mostrar el diálogo ANTES de iniciar el procesamiento
        progress_dialog.show()
        QApplication.processEvents() # Procesar eventos para que se muestre
        
        # Crear e iniciar el worker DESPUÉS de mostrar el diálogo
        self.worker = GenerateImageWorker(modelo_imagen, titulo, contenido, imagen_personalizada)
        
        # Conectar señales
        self.worker.signals.success.connect(
            lambda path: self.handle_image_success(path, progress_dialog))
        self.worker.signals.error.connect(
            lambda error: self.handle_image_error(error, progress_dialog))
        self.worker.signals.finished.connect(
            lambda: self.worker.deleteLater())  # Limpiar el worker después
        # --- NUEVO: Conectar señal de cancelación --- 
        progress_dialog.canceled.connect(self.cancel_image_generation)
        
        # Iniciar el worker
        self.worker.start()
    
    def handle_image_success(self, image_path, progress_dialog):
        """Manejar el éxito de la generación de la imagen"""
        progress_dialog.close()
        
        # Actualizar la ruta de la imagen y la etiqueta
        self.new_image_path = image_path
        self.current_image_label.setText(os.path.basename(image_path))
        
        # Actualizar la vista previa temporalmente
        self.update_preview_temporarily()
        
        # Mostrar mensaje de éxito
        QMessageBox.information(self,
                                obtener_traduccion('ai_image_success_title', self.current_language),
                                obtener_traduccion('ai_image_success_message', self.current_language))
    
    def handle_image_error(self, error, progress_dialog):
        """Manejar el error de la generación de la imagen"""
        progress_dialog.close()
        error_message = obtener_traduccion('ai_image_error_message', self.current_language).format(error=error)
        QMessageBox.critical(self, obtener_traduccion('error', self.current_language), error_message)
        
    def update_preview_temporarily(self):
        """Actualizar la vista previa temporalmente sin guardar los cambios"""
        if not self.parent_widget or not isinstance(self.parent_widget, VentanaVistaPrevia):
            return
            
        self.temp_preview_updated = True
        
        # Guardar los datos originales de la diapositiva actual (solo la imagen si es lo único que cambia aquí)
        # self.original_image_path se maneja en browse_new_image y generación IA
        
        # Aplicar los cambios de formato (current_format_settings del diálogo) TEMPORALMENTE a VentanaVistaPrevia
        # para que la vista previa refleje estos cambios.
        # Guardamos el estado global de VentanaVistaPrevia para poder restaurarlo.
        # Esto se hace en load_format_settings y se guarda en self.original_format_settings del diálogo.
        
        # Clonamos original_format_settings para no modificarlo por referencia
        temp_parent_original_format = dict(self.original_format_settings)
        
        # Aplicamos los cambios actuales del diálogo a la VentanaVistaPrevia para la vista previa
        self.parent_widget.update_format_settings(
            self.current_format_settings['title_font_name'],
            self.current_format_settings['content_font_name'],
            self.current_format_settings['title_font_size'],
            self.current_format_settings['content_font_size'],
            self.current_format_settings['title_bold'],
            self.current_format_settings['title_italic'],
            self.current_format_settings['title_underline'],
            self.current_format_settings['content_bold'],
            self.current_format_settings['content_italic'],
            self.current_format_settings['content_underline']
        )
        
        # Crear datos temporales para la diapositiva actual (título, contenido, imagen)
        temp_slide_data = {
            'imagen_path': self.new_image_path, # self.new_image_path se actualiza en browse_new_image o IA
            'titulo': self.title_edit.text(),
            'contenido': self.content_edit.toPlainText(),
            # No guardamos format_settings aquí, porque la vista previa ya usa el formato global
            # de VentanaVistaPrevia, que acabamos de actualizar temporalmente.
            'format_settings': self.current_format_settings  # Incluir la configuración específica de formato
        }
        
        # Guardar temporalmente los nuevos datos de la diapositiva
        current_slide_idx = self.parent_widget.current_slide_index
        self.parent_widget.all_slides_data[current_slide_idx] = temp_slide_data
        
        # Actualizar la vista previa
        self.parent_widget.mostrar_diapositiva_actual()
        
        # Ya no es necesario restaurar el formato global porque ahora cada diapositiva usa su propio formato
        # self.parent_widget.update_format_settings(**temp_parent_original_format)
        
    def on_dialog_finished(self, result):
        """Manejar el cierre del diálogo, restaurando vista previa si es necesario"""
        if self.parent_widget and isinstance(self.parent_widget, VentanaVistaPrevia):
            if result == QDialog.Rejected and self.temp_preview_updated:
                # Si se canceló y hubo cambios temporales, restaurar datos y formato original de la diapositiva
                current_slide_idx = self.parent_widget.current_slide_index
                self.parent_widget.all_slides_data[current_slide_idx] = dict(self.initial_data)
                
                # Restaurar la configuración de formato global de VentanaVistaPrevia a como estaba
                # ANTES de que se abriera este diálogo de edición para esta diapositiva.
                # self.original_format_settings del diálogo guarda ese estado.
                self.parent_widget.update_format_settings(**self.original_format_settings)
                
                self.parent_widget.mostrar_diapositiva_actual()
            elif result == QDialog.Accepted:
                # Si se aceptó, la configuración de formato actual del diálogo (current_format_settings)
                # se convertirá en la configuración específica de esta diapositiva.
                # Y también se aplicará como la nueva configuración global de VentanaVistaPrevia.
                current_slide_idx = self.parent_widget.current_slide_index
                edited_data = self.get_edited_data() # Esto ya incluye format_settings
                self.parent_widget.all_slides_data[current_slide_idx] = edited_data
                
                # Aplicar la configuración de formato del diálogo a VentanaVistaPrevia como la nueva global
                self.parent_widget.update_format_settings(**self.current_format_settings)
                self.parent_widget.mostrar_diapositiva_actual() # Refrescar para mostrar con el formato final

    def get_edited_data(self):
        # Incluir también la configuración de formato en los datos editados
        return {
            'titulo': self.title_edit.text(),
            'contenido': self.content_edit.toPlainText(),
            'imagen_path': self.new_image_path,
            'format_settings': self.current_format_settings
        }

    # --- NUEVA FUNCIÓN para manejar la cancelación ---
    def cancel_image_generation(self):
        print("Señal de cancelación recibida en EditSlideDialog.")
        if hasattr(self, 'worker') and self.worker and self.worker.isRunning():
            print("Solicitando interrupción del worker desde EditSlideDialog...")
            self.worker.requestInterruption()
        else:
            print("Worker no encontrado o no está corriendo en EditSlideDialog.")
        # El diálogo de progreso se cierra automáticamente al cancelar (no es necesario cerrarlo aquí)

    def accept(self):
        # Ya no se llama a list_widget aquí, se asume que este método es solo para QDialogButtonBox.Save
        # La lógica de actualizar VentanaVistaPrevia con current_format_settings se mueve a on_dialog_finished
        super().accept()
# --- Fin Clase EditSlideDialog ---

# Clase para la ventana de vista previa de diapositivas
class VentanaVistaPrevia(QWidget):
    # --- MODIFICADO: Añadir parámetro initial_language y eliminar la lógica de búsqueda en padre/abuelo ---
    def __init__(self, parent=None, initial_language='es', # Añadir parámetro
                 title_font_name='Calibri', content_font_name='Calibri', 
                 title_font_size=16, content_font_size=10, 
                 title_bold=False, title_italic=False, title_underline=False, 
                 content_bold=False, content_italic=False, content_underline=False):
        super().__init__(parent)
        # --- MODIFICADO: Usar initial_language directamente ---
        self.current_language = initial_language
        # --- ELIMINADO: Lógica de búsqueda en padre/abuelo ---
        # self.current_language = 'es' # Default inicial
        # # Obtener el idioma del padre si está disponible
        # if parent and hasattr(parent, 'current_language'):
        #     self.current_language = parent.current_language
        # elif parent and hasattr(parent, 'parent') and parent.parent and hasattr(parent.parent, 'current_language'):
        #      # Intentar obtener del "abuelo" si el padre directo no lo tiene (caso común en anidamiento)
        #     self.current_language = parent.parent.current_language
        # -------------------------------------------------------

        self.pptx_path = None # Ruta al archivo PPTX generado
        
        # --- NUEVO: Almacenar configuración de formato global ---
        self.title_font_name = title_font_name
        self.content_font_name = content_font_name
        self.title_font_size = title_font_size
        self.content_font_size = content_font_size
        self.title_bold = title_bold
        self.title_italic = title_italic
        self.title_underline = title_underline
        self.content_bold = content_bold
        self.content_italic = content_italic
        self.content_underline = content_underline
        # -------------------------------------------------------
        
        self.setMinimumSize(600, 350)  # Reducir altura mínima
        
        # Inicializar la lista de todas las diapositivas
        self.all_slides_data = []
        self.current_slide_index = -1
        
        # Variable para controlar si el usuario está navegando manualmente
        self.modo_navegacion = False
        
        # Crear layout principal
        layout = QVBoxLayout()
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(5)  # Reducir espacio
        
        # Crear área de desplazamiento
        self.scroll = QScrollArea()
        self.scroll.setWidgetResizable(True)
        self.scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.scroll.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.scroll.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        # Widget contenedor para la diapositiva actual
        self.contenedor = QWidget()
        self.contenedor.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.contenedor_layout = QVBoxLayout()
        self.contenedor_layout.setAlignment(Qt.AlignCenter)
        self.contenedor_layout.setContentsMargins(5, 5, 5, 5)  # Reducir márgenes
        self.contenedor_layout.setSpacing(5)  # Reducir espaciado
        
        # Añadir mensaje informativo cuando no hay diapositivas
        self.mensaje_vacio = QLabel()
        self.mensaje_vacio.setAlignment(Qt.AlignCenter)
        self.mensaje_vacio.setFont(QFont("Arial", 12))  # Reducir tamaño de fuente
        self.mensaje_vacio.setStyleSheet("color: #888; padding: 50px;")  # Reducir padding
        # --- MODIFICADO: Establecer texto inicial usando el idioma determinado ---
        mensaje_inicial = obtener_traduccion('vista_previa_empty_message', self.current_language)
        self.mensaje_vacio.setText(mensaje_inicial if mensaje_inicial != 'vista_previa_empty_message' else "Aquí aparecerán las diapositivas generadas")
        # ---------------------------------------------------------------------
        self.mensaje_vacio.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        # Añadir imagen de previsualización
        self.imagen_preview = QLabel()
        self.imagen_preview.setAlignment(Qt.AlignCenter)
        pixmap = QPixmap(resource_path("iconos/icon.png"))
        if not pixmap.isNull():
            pixmap = pixmap.scaled(128, 128, Qt.KeepAspectRatio, Qt.SmoothTransformation)
            self.imagen_preview.setPixmap(pixmap)
        self.imagen_preview.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        preview_container = QVBoxLayout()
        preview_container.setAlignment(Qt.AlignCenter)
        preview_container.addWidget(self.imagen_preview)
        preview_container.addWidget(self.mensaje_vacio)
        
        # Contenedor para el mensaje vacío y la imagen
        empty_widget = QWidget()
        empty_widget.setLayout(preview_container)
        empty_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        
        self.contenedor_layout.addWidget(empty_widget, 1, Qt.AlignCenter)
        
        self.contenedor.setLayout(self.contenedor_layout)
        
        self.scroll.setWidget(self.contenedor)
        layout.addWidget(self.scroll, 1)
        
        # Añadir botones de navegación
        nav_layout = QHBoxLayout()
        nav_layout.setContentsMargins(5, 0, 5, 0)  # Reducir márgenes
        nav_layout.setSpacing(10)  # Reducir espacio
        
        # Inicializar los botones con texto traducido
        # --- MODIFICADO: Usar self.current_language determinado al inicio ---
        # texto_anterior = obtener_traduccion('slide_previous', self.current_language)
        # self.prev_button = QPushButton(f"< {texto_anterior}")
        # # ------------------------------------------------------------------
        # self.prev_button.setEnabled(False)
        # self.prev_button.clicked.connect(self.mostrar_diapositiva_anterior)
        # self.prev_button.setSizePolicy(QSizePolicy.Preferred, QSizePolicy.Fixed)
        # self.prev_button.setFixedHeight(25)  # Altura reducida
        
        # --- NUEVO: Botón Anterior con Icono ---
        self.prev_button = QPushButton()
        left_icon_path = resource_path("iconos/left.png")
        if os.path.exists(left_icon_path):
            self.prev_button.setIcon(QIcon(left_icon_path))
        else:
            self.prev_button.setText("<") # Fallback
        self.prev_button.setIconSize(QSize(32, 32))
        self.prev_button.setFixedSize(QSize(40, 40))
        self.prev_button.setToolTip(obtener_traduccion('slide_previous', self.current_language))
        self.prev_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
        self.prev_button.setCursor(Qt.PointingHandCursor)
        self.prev_button.setEnabled(False)
        self.prev_button.clicked.connect(self.mostrar_diapositiva_anterior)
        # --- FIN NUEVO ---

        self.slide_counter_label = QLabel("0/0")
        self.slide_counter_label.setAlignment(Qt.AlignCenter)
        self.slide_counter_label.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        
        # --- NUEVO: Aumentar tamaño de fuente para el contador ---
        font_contador = QFont()
        font_contador.setPointSize(17) # Tamaño de fuente más grande
        self.slide_counter_label.setFont(font_contador)
        # --- FIN NUEVO ---

        # --- NUEVO: Botón Siguiente con Icono ---
        self.next_button = QPushButton()
        right_icon_path = resource_path("iconos/right.png")
        if os.path.exists(right_icon_path):
            self.next_button.setIcon(QIcon(right_icon_path))
        else:
            self.next_button.setText(">") # Fallback
        self.next_button.setIconSize(QSize(32, 32))
        self.next_button.setFixedSize(QSize(40, 40))
        self.next_button.setToolTip(obtener_traduccion('slide_next', self.current_language))
        self.next_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
        self.next_button.setCursor(Qt.PointingHandCursor)
        self.next_button.setEnabled(False)
        self.next_button.clicked.connect(self.mostrar_diapositiva_siguiente)
        # --- FIN NUEVO ---
        
        nav_layout.addWidget(self.prev_button)
        nav_layout.addWidget(self.slide_counter_label)
        nav_layout.addWidget(self.next_button)
        
        layout.addLayout(nav_layout)
        
        self.setLayout(layout)
        
        # Guardar las proporciones originales para el escalado responsivo
        self.original_width = self.width()
        self.original_height = self.height()
        
        # Conectar evento de redimensionamiento
        self.resized = False
        
        # Aplicar traducciones al iniciar - asegurarnos de que se actualicen todos los textos
        # --- MODIFICADO: Llamar a actualizar_idioma al final para asegurar la consistencia ---
        # Esta llamada ahora reforzará los textos ya establecidos con el idioma correcto
        self.actualizar_idioma(self.current_language)
        # ------------------------------------------------------------------------------------
        
        # Habilitar eventos de rueda del mouse para navegación entre diapositivas
        self.setMouseTracking(True)
        
        # Variable para almacenar el botón de edición (se crea/destruye dinámicamente)
        self.edit_button = None
        # Variable para almacenar el botón de abrir (se crea/destruye dinámicamente)
        self.open_button = None
    
    # --- NUEVO: Método para actualizar la configuración de formato ---
    def update_format_settings(self, title_font_name, content_font_name, 
                               title_font_size, content_font_size, 
                               title_bold, title_italic, title_underline, 
                               content_bold, content_italic, content_underline):
        self.title_font_name = title_font_name
        self.content_font_name = content_font_name
        self.title_font_size = title_font_size
        self.content_font_size = content_font_size
        self.title_bold = title_bold
        self.title_italic = title_italic
        self.title_underline = title_underline
        self.content_bold = content_bold
        self.content_italic = content_italic
        self.content_underline = content_underline
        print("Configuración de formato de VistaPrevia actualizada.") # Debug
        # Nota: Ya no actualizamos la diapositiva automáticamente aquí
    # -------------------------------------------------------------

    # --- Nueva función para establecer la ruta del PPTX ---
    def set_pptx_path(self, path):
        if os.path.exists(path):
            self.pptx_path = path
        else:
            self.pptx_path = None
            print(f"Advertencia: La ruta del PPTX proporcionada no existe: {path}")
            
    # Manejar eventos de rueda del mouse para la navegación entre diapositivas
    def wheelEvent(self, event):
        # Solo procesar el evento si hay diapositivas
        if self.all_slides_data and len(self.all_slides_data) > 0:
            # Detectar dirección del scroll
            if event.angleDelta().y() > 0:
                # Scroll hacia arriba - ir a diapositiva anterior
                self.mostrar_diapositiva_anterior()
            else:
                # Scroll hacia abajo - ir a diapositiva siguiente
                self.mostrar_diapositiva_siguiente()
            
            # Consumir el evento para evitar comportamiento estándar
            event.accept()
        else:
            # Si no hay diapositivas, permitir el comportamiento normal de scroll
            super().wheelEvent(event)
    
    # Capturar evento de redimensionamiento
    def resizeEvent(self, event: QResizeEvent):
        super().resizeEvent(event)
        if self.current_slide_index >= 0 and not self.resized:
            # Usar un temporizador para evitar múltiples actualizaciones durante el redimensionamiento
            self.resized = True
            self.mostrar_diapositiva_actual()
            self.resized = False
    
    # Función para actualizar el idioma de la interfaz
    def actualizar_idioma(self, language):
        try:
            # Establecer el idioma actual
            self.current_language = language
            
            # Actualizar título de la ventana de forma segura
            try:
                self.setWindowTitle(obtener_traduccion('vista_previa_title', self.current_language))
            except:
                pass
            
            # Actualizar textos de forma segura sin intentar modificar widgets que puedan estar eliminados
            try:
                # Actualizar solo los textos de los botones
                self.actualizar_textos_botones()
            except Exception as e:
                print(f"Error al actualizar botones: {str(e)}")
            
            # Si no hay diapositivas, actualizar el mensaje vacío
            try:
                mensaje = obtener_traduccion('vista_previa_empty_message', self.current_language)
                if mensaje:
                    if hasattr(self, 'mensaje_vacio') and self.mensaje_vacio:
                        self.mensaje_vacio.setText(mensaje)
            except Exception as e:
                print(f"Error al actualizar mensaje vacío: {str(e)}")
            
            # Si no hay diapositivas, forzar una reconstrucción completa de la vista
            if not self.all_slides_data or self.current_slide_index < 0:
                try:
                    # Eliminar completamente el contenido actual y recrearlo
                    self.limpiar_contenedor()
                    
                    # Recrear la vista desde cero
                    self.mostrar_diapositiva_actual()
                except Exception as e:
                    print(f"Error al recrear vista vacía: {str(e)}")
            else:
                # Si hay diapositivas, simplemente actualizar la vista actual
                try:
                    self.mostrar_diapositiva_actual()
                    # Actualizar texto del botón editar si existe
                    if self.edit_button:
                         self.edit_button.setToolTip(obtener_traduccion('edit_slide', self.current_language))
                    # Actualizar texto del botón abrir si existe
                    if self.open_button:
                         self.open_button.setToolTip(obtener_traduccion('open_slide', self.current_language))
                except Exception as e:
                    print(f"Error al actualizar diapositiva actual: {str(e)}")
        except Exception as e:
            print(f"Error en actualizar_idioma: {str(e)}")
    
    # Función para actualizar textos de los botones de navegación
    def actualizar_textos_botones(self):
        try:
            # Obtener traducciones
            texto_anterior = obtener_traduccion('slide_previous', self.current_language)
            if not texto_anterior:
                texto_anterior = "Anterior"
                
            texto_siguiente = obtener_traduccion('slide_next', self.current_language)
            if not texto_siguiente:
                texto_siguiente = "Siguiente"
            
            # Verificar si los botones existen antes de actualizarlos
            if hasattr(self, 'prev_button') and self.prev_button is not None:
                try:
                    # self.prev_button.setText("< " + texto_anterior) # Ya no se establece texto
                    self.prev_button.setToolTip(texto_anterior) # Actualizar tooltip
                except RuntimeError:
                    # Si el botón ya fue eliminado, recrearlo con el nuevo estilo
                    # self.prev_button = QPushButton("< " + texto_anterior)
                    self.prev_button = QPushButton()
                    left_icon_path = resource_path("iconos/left.png")
                    if os.path.exists(left_icon_path):
                        self.prev_button.setIcon(QIcon(left_icon_path))
                    else:
                        self.prev_button.setText("<")
                    self.prev_button.setIconSize(QSize(32, 32))
                    self.prev_button.setFixedSize(QSize(40, 40))
                    self.prev_button.setToolTip(texto_anterior)
                    self.prev_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
                    self.prev_button.setCursor(Qt.PointingHandCursor)
                    # El estado de enabled se maneja en actualizar_botones_navegacion,
                    # pero mantenemos el connect y una configuración inicial segura.
                    self.prev_button.setEnabled(False) 
                    self.prev_button.clicked.connect(self.mostrar_diapositiva_anterior)
                    # self.prev_button.setSizePolicy(QSizePolicy.Preferred, QSizePolicy.Fixed) # Ya no es necesario
                    # self.prev_button.setFixedHeight(25) # Ya no es necesario
                
            if hasattr(self, 'next_button') and self.next_button is not None:
                try:
                    # self.next_button.setText(texto_siguiente + " >") # Ya no se establece texto
                    self.next_button.setToolTip(texto_siguiente) # Actualizar tooltip
                except RuntimeError:
                    # Si el botón ya fue eliminado, recrearlo con el nuevo estilo
                    # self.next_button = QPushButton(texto_siguiente + " >")
                    self.next_button = QPushButton()
                    right_icon_path = resource_path("iconos/right.png")
                    if os.path.exists(right_icon_path):
                        self.next_button.setIcon(QIcon(right_icon_path))
                    else:
                        self.next_button.setText(">")
                    self.next_button.setIconSize(QSize(32, 32))
                    self.next_button.setFixedSize(QSize(40, 40))
                    self.next_button.setToolTip(texto_siguiente)
                    self.next_button.setStyleSheet("QPushButton { border: none; background-color: transparent; }")
                    self.next_button.setCursor(Qt.PointingHandCursor)
                    self.next_button.setEnabled(False)
                    self.next_button.clicked.connect(self.mostrar_diapositiva_siguiente)
                    # self.next_button.setSizePolicy(QSizePolicy.Preferred, QSizePolicy.Fixed) # Ya no es necesario
                    # self.next_button.setFixedHeight(25) # Ya no es necesario
        except Exception as e:
            print(f"Error al actualizar textos de botones: {str(e)}")
    
    # Función para mostrar la diapositiva actual
    def mostrar_diapositiva_actual(self):
        try:
            # Limpiar el contenedor de diapositivas
            self.limpiar_contenedor()
            self.edit_button = None # Resetear referencia al botón
            self.open_button = None # Resetear referencia al botón abrir
            
            # Verificar que hay diapositivas y que el índice es válido
            if not self.all_slides_data or self.current_slide_index < 0:
                # Recrear el mensaje vacío para asegurar que no haya problemas con objetos eliminados
                empty_widget = QWidget()
                preview_container = QVBoxLayout(empty_widget)
                preview_container.setAlignment(Qt.AlignCenter)
                
                # Recrear la imagen de previsualización
                imagen_preview = QLabel()
                imagen_preview.setAlignment(Qt.AlignCenter)
                pixmap = QPixmap(resource_path("iconos/icon.png"))
                if not pixmap.isNull():
                    pixmap = pixmap.scaled(128, 128, Qt.KeepAspectRatio, Qt.SmoothTransformation)
                    imagen_preview.setPixmap(pixmap)
                imagen_preview.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
                
                # Recrear el mensaje vacío - IMPORTANTE: siempre usar el idioma actual
                mensaje_vacio = QLabel()
                mensaje_vacio.setAlignment(Qt.AlignCenter)
                mensaje_vacio.setFont(QFont("Arial", 12))
                mensaje_vacio.setStyleSheet("color: #888; padding: 50px;")
                mensaje_vacio.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
                
                # Obtener el texto traducido con manejo de errores
                try:
                    mensaje = obtener_traduccion('vista_previa_empty_message', self.current_language)
                    if not mensaje:
                        mensaje = "Aquí aparecerán las diapositivas generadas"
                except Exception:
                    mensaje = "Aquí aparecerán las diapositivas generadas"
                
                mensaje_vacio.setText(mensaje)
                
                # Guardar referencia al mensaje vacío nuevo
                self.mensaje_vacio = mensaje_vacio
                
                preview_container.addWidget(imagen_preview)
                preview_container.addWidget(mensaje_vacio)
                self.contenedor_layout.addWidget(empty_widget, 1, Qt.AlignCenter)
                return
            
            # Obtener datos de la diapositiva actual
            datos = self.all_slides_data[self.current_slide_index]
            
            # Obtener la configuración de formato específica de esta diapositiva si existe
            # Si no existe, usar la configuración global actual
            format_settings = datos.get('format_settings', {})
            
            # Usar formato específico de la diapositiva o la configuración global como fallback
            title_font_name = format_settings.get('title_font_name', self.title_font_name)
            content_font_name = format_settings.get('content_font_name', self.content_font_name)
            title_font_size = format_settings.get('title_font_size', self.title_font_size)
            content_font_size = format_settings.get('content_font_size', self.content_font_size)
            title_bold = format_settings.get('title_bold', self.title_bold)
            title_italic = format_settings.get('title_italic', self.title_italic)
            title_underline = format_settings.get('title_underline', self.title_underline)
            content_bold = format_settings.get('content_bold', self.content_bold)
            content_italic = format_settings.get('content_italic', self.content_italic)
            content_underline = format_settings.get('content_underline', self.content_underline)
            
            # Crear frame para la diapositiva
            slide_frame = QFrame()
            # slide_frame.setFrameStyle(QFrame.Box) # Opcional: quitar si no se desea borde
            slide_frame.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
            slide_layout = QVBoxLayout()
            slide_layout.setSpacing(8) # Mantener espaciado actual o ajustar
            slide_layout.setContentsMargins(8, 8, 8, 8) # Mantener márgenes actuales o ajustar
            
            # --- Layout Superior (Título y Botones) ---
            top_layout = QHBoxLayout()
            top_layout.setContentsMargins(0,0,0,0)
            top_layout.setSpacing(5)
            
            # --- Crear Botón Abrir ---
            self.open_button = QPushButton()
            icon_path_open = resource_path("iconos/open.png")
            if os.path.exists(icon_path_open):
                self.open_button.setIcon(QIcon(icon_path_open))
                self.open_button.setIconSize(QSize(32, 32))
                self.open_button.setFixedSize(QSize(40, 40))
            else:
                self.open_button.setText("O") # Texto fallback si no hay icono
                self.open_button.setFixedSize(QSize(40, 40))
                
            open_tooltip = obtener_traduccion('open_slide', self.current_language)
            self.open_button.setToolTip(open_tooltip if open_tooltip != 'open_slide' else "Abrir diapositiva")
            self.open_button.setStyleSheet("QPushButton { border: none; background-color: transparent; padding: 2px; }")
            self.open_button.setCursor(Qt.PointingHandCursor)
            self.open_button.clicked.connect(self.abrir_presentacion)
            self.open_button.setEnabled(self.pptx_path is not None)
            
            titulo_label = QLabel(datos['titulo'])
            titulo_label.setAlignment(Qt.AlignCenter)
            titulo_label.setWordWrap(True)
            
            final_title_size = title_font_size 
            title_font_obj = QFont(title_font_name, final_title_size)
            title_font_obj.setBold(title_bold)
            title_font_obj.setItalic(title_italic)
            title_font_obj.setUnderline(title_underline)
            titulo_label.setFont(title_font_obj)
            
            self.edit_button = QPushButton()
            icon_path = resource_path("iconos/editar.png")
            if os.path.exists(icon_path):
                self.edit_button.setIcon(QIcon(icon_path))
                self.edit_button.setIconSize(QSize(32, 32))
                self.edit_button.setFixedSize(QSize(40, 40))
            else:
                self.edit_button.setText("E") 
                self.edit_button.setFixedSize(QSize(40, 40))
                
            edit_tooltip = obtener_traduccion('edit_slide', self.current_language)
            self.edit_button.setToolTip(edit_tooltip if edit_tooltip != 'edit_slide' else "Editar diapositiva")
            self.edit_button.setStyleSheet("QPushButton { border: none; background-color: transparent; padding: 2px; }")
            self.edit_button.setCursor(Qt.PointingHandCursor)
            self.edit_button.clicked.connect(self.editar_diapositiva_actual)
            self.edit_button.setEnabled(PPTX_AVAILABLE and self.pptx_path is not None) 
            
            top_layout.addWidget(self.open_button, 0, Qt.AlignLeft | Qt.AlignTop)
            top_layout.addWidget(titulo_label, 1) 
            top_layout.addWidget(self.edit_button, 0, Qt.AlignRight | Qt.AlignTop)
            
            # --- Contenido (crear y calcular su altura necesaria ANTES de la imagen) ---
            contenido_label = QLabel()
            contenido_label.setTextFormat(Qt.PlainText) 
            contenido_label.setText(datos['contenido']) 
            contenido_label.setAlignment(Qt.AlignLeft) # o Qt.AlignJustify
            contenido_label.setWordWrap(True)
            
            final_content_size = content_font_size
            content_font_obj = QFont(content_font_name, final_content_size)
            content_font_obj.setBold(content_bold)
            content_font_obj.setItalic(content_italic)
            content_font_obj.setUnderline(content_underline)
            contenido_label.setFont(content_font_obj)

            # Calcular el ancho efectivo para el texto (título y contenido)
            # Esta es la anchura disponible DENTRO de slide_layout para sus hijos.
            effective_text_width = (self.scroll.viewport().width() -
                                   (self.contenedor_layout.contentsMargins().left() + self.contenedor_layout.contentsMargins().right()) -
                                   (slide_frame.contentsMargins().left() + slide_frame.contentsMargins().right()) -
                                   (slide_layout.contentsMargins().left() + slide_layout.contentsMargins().right()) -
                                   20) # Buffer adicional (e.g., para posible scrollbar si el cálculo no es perfecto)
            if effective_text_width <= 20: # Fallback si el cálculo da un ancho muy pequeño
                effective_text_width = int(self.scroll.viewport().width() * 0.85)


            # Calcular la altura necesaria para el contenido_label:
            fm_content = QFontMetrics(content_font_obj)
            # boundingRect para calcular la altura del texto con word wrap.
            text_rect_content = fm_content.boundingRect(QRect(0, 0, int(effective_text_width), 9999),
                                                       int(Qt.TextWordWrap | Qt.AlignLeft), datos['contenido'])
            needed_content_height = text_rect_content.height() + fm_content.leading() // 2 # Pequeño buffer (ej. medio interlineado)
            min_content_height = fm_content.height() # Altura mínima de una línea
            needed_content_height = max(needed_content_height, min_content_height)
            contenido_label.setMaximumHeight(needed_content_height) # FIJAR SU ALTURA MÁXIMA

            # Estimar altura del top_layout (título + botones)
            fm_title = QFontMetrics(title_font_obj)
            title_rect = fm_title.boundingRect(QRect(0, 0, int(effective_text_width), 9999),
                                              int(Qt.TextWordWrap | Qt.AlignCenter), datos['titulo'])
            estimated_title_text_height = title_rect.height()
            
            button_height_estimate = self.edit_button.sizeHint().height() # Usar sizeHint del botón
            if button_height_estimate <=0: button_height_estimate = 32 # Fallback
            
            top_layout_height_estimate = (max(estimated_title_text_height, button_height_estimate) +
                                          top_layout.contentsMargins().top() + top_layout.contentsMargins().bottom())


            # Calcular altura disponible para la imagen
            scroll_viewport_h = self.scroll.viewport().height()

            # Márgenes verticales totales desde el viewport hasta el interior de slide_layout
            total_v_margins_around_slide_layout_children = (
                self.contenedor_layout.contentsMargins().top() + self.contenedor_layout.contentsMargins().bottom() +
                slide_frame.contentsMargins().top() + slide_frame.contentsMargins().bottom() +
                slide_layout.contentsMargins().top() + slide_layout.contentsMargins().bottom()
            )
            
            # Espaciados dentro del slide_layout (entre top_layout, imagen, contenido)
            num_gaps_in_slide_layout = 2 
            slide_layout_total_spacing = slide_layout.spacing() * num_gaps_in_slide_layout

            total_fixed_height_plus_margins_spacing = (top_layout_height_estimate + needed_content_height +
                                                       total_v_margins_around_slide_layout_children +
                                                       slide_layout_total_spacing)

            available_height_for_image = scroll_viewport_h - total_fixed_height_plus_margins_spacing
            min_image_render_height = 30 # Mínimo absoluto en píxeles para que la imagen no desaparezca
            available_height_for_image = max(available_height_for_image, min_image_render_height)

            # --- Imagen (crear y escalar) ---
            imagen_label = QLabel()
            imagen_label.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding) # Importante
            imagen_label.setAlignment(Qt.AlignCenter)

            pixmap = QPixmap(datos['imagen_path'])
            if not pixmap.isNull():
                # Usar effective_text_width para el ancho máximo de la imagen, ya que es el ancho disponible.
                max_img_width_px = int(effective_text_width)
                
                scaled_pixmap = pixmap.scaled(
                    max_img_width_px,
                    int(available_height_for_image), # Usar la altura calculada
                    Qt.KeepAspectRatio,
                    Qt.SmoothTransformation
                )
                imagen_label.setPixmap(scaled_pixmap)
            else:
                 imagen_label.setText(f"Error al cargar: {os.path.basename(datos['imagen_path'])}")
                 imagen_label.setStyleSheet("color: red;")
            
            # --- Añadir al layout en orden ---
            slide_layout.addLayout(top_layout)
            slide_layout.addWidget(imagen_label) 
            slide_layout.addWidget(contenido_label)
            
            slide_frame.setLayout(slide_layout)
            self.contenedor_layout.addWidget(slide_frame)
            
            # Actualizar el contador de diapositivas
            self.slide_counter_label.setText(f"{self.current_slide_index + 1}/{len(self.all_slides_data)}")
            
            # Si estamos viendo la última diapositiva, desactivar el modo navegación
            if self.current_slide_index == len(self.all_slides_data) - 1:
                self.modo_navegacion = False
            
            # Asegurar que la diapositiva sea visible y se ajuste
            self.contenedor.adjustSize() # Puede ayudar a que el layout se recalcule
        except Exception as e:
            print(f"Error en mostrar_diapositiva_actual: {str(e)}")
            import traceback
            traceback.print_exc()

    # Función para limpiar el contenedor de diapositivas
    def limpiar_contenedor(self):
        # Eliminar todos los widgets del contenedor
        while self.contenedor_layout.count():
            item = self.contenedor_layout.takeAt(0)
            if item.widget():
                item.widget().deleteLater()

    # Función para mostrar la diapositiva anterior
    def mostrar_diapositiva_anterior(self):
        if self.current_slide_index > 0:
            # Activar el modo navegación
            self.modo_navegacion = True
            self.current_slide_index -= 1
            self.mostrar_diapositiva_actual()
            self.actualizar_botones_navegacion()

    # Función para mostrar la diapositiva siguiente
    def mostrar_diapositiva_siguiente(self):
        if self.current_slide_index < len(self.all_slides_data) - 1:
            # Activar el modo navegación si no estamos en la última
            if self.current_slide_index < len(self.all_slides_data) - 2:
                self.modo_navegacion = True
            else:
                # Si vamos a la última, desactivar modo navegación
                self.modo_navegacion = False
                
            self.current_slide_index += 1
            self.mostrar_diapositiva_actual()
            self.actualizar_botones_navegacion()

    # Función para actualizar el estado de los botones de navegación
    def actualizar_botones_navegacion(self):
        # Activar/desactivar botones según la posición actual
        self.prev_button.setEnabled(self.current_slide_index > 0)
        self.next_button.setEnabled(self.current_slide_index < len(self.all_slides_data) - 1) 
        
    # Función para resetear completamente la vista previa
    def reset_completo(self):
        try:
            # Limpiar la lista de diapositivas
            self.all_slides_data = []
            # Reiniciar el índice
            self.current_slide_index = -1
            # Desactivar botones de navegación
            self.prev_button.setEnabled(False)
            self.next_button.setEnabled(False)
            # Resetear contador
            self.slide_counter_label.setText("0/0")
            # Desactivar modo navegación
            self.modo_navegacion = False
            # Resetear botón editar
            self.edit_button = None 
            # Resetear botón abrir
            self.open_button = None
            
            # Limpiar el contenedor y mostrar mensaje vacío
            self.limpiar_contenedor()
            
            # Recrear el mensaje vacío para asegurar que no haya problemas con objetos eliminados
            empty_widget = QWidget()
            empty_widget.setStyleSheet("background-color: transparent;")
            preview_container = QVBoxLayout(empty_widget)
            preview_container.setAlignment(Qt.AlignCenter)
            preview_container.setContentsMargins(0, 0, 0, 0)
            preview_container.setSpacing(10)
            
            # Recrear la imagen de previsualización
            imagen_preview = QLabel()
            imagen_preview.setAlignment(Qt.AlignCenter)
            pixmap = QPixmap(resource_path("iconos/icon.png"))
            if not pixmap.isNull():
                pixmap = pixmap.scaled(128, 128, Qt.KeepAspectRatio, Qt.SmoothTransformation)
                imagen_preview.setPixmap(pixmap)
            imagen_preview.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
            imagen_preview.setStyleSheet("background-color: transparent;")
            
            # Recrear el mensaje vacío - IMPORTANTE: siempre usar el idioma actual
            mensaje_vacio = QLabel()
            mensaje_vacio.setAlignment(Qt.AlignCenter)
            mensaje_vacio.setFont(QFont("Arial", 12))
            mensaje_vacio.setStyleSheet("color: #888; padding: 50px; background-color: transparent;")
            mensaje_vacio.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
            
            # Obtener el texto traducido con manejo de errores
            try:
                mensaje = obtener_traduccion('vista_previa_empty_message', self.current_language)
                if not mensaje:
                    mensaje = "Aquí aparecerán las diapositivas generadas"
            except Exception:
                mensaje = "Aquí aparecerán las diapositivas generadas"
            
            mensaje_vacio.setText(mensaje)
            
            # Guardar referencia al nuevo mensaje vacío
            self.mensaje_vacio = mensaje_vacio
            
            preview_container.addWidget(imagen_preview)
            preview_container.addWidget(mensaje_vacio)
            
            # Limpiar el layout anterior si existe
            while self.contenedor_layout.count():
                item = self.contenedor_layout.takeAt(0)
                if item.widget():
                    item.widget().deleteLater()
            
            # Agregar el nuevo widget al contenedor
            self.contenedor_layout.addWidget(empty_widget, 1, Qt.AlignCenter)
            
            # Forzar actualización visual
            empty_widget.update()
            self.update()
            QApplication.processEvents()
            
            self.pptx_path = None # Olvidar la ruta del PPTX al resetear
            
        except Exception as e:
            print(f"Error en reset_completo: {str(e)}")

    # --- NUEVA FUNCIÓN: Abrir presentación en el programa predeterminado del sistema ---
    def abrir_presentacion(self):
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            QMessageBox.warning(self,
                                obtener_traduccion('open_pptx_error_title', self.current_language),
                                obtener_traduccion('open_pptx_not_found_message', self.current_language))
            return
            
        try:
            # En Windows usar 'start', en otros sistemas puede variar
            if sys.platform == 'win32':
                os.startfile(self.pptx_path)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', self.pptx_path])
            else:  # Linux y otros
                subprocess.call(['xdg-open', self.pptx_path])
                
            print(f"Presentación abierta: {self.pptx_path}")
        except Exception as e:
            print(f"Error al abrir la presentación: {e}")
            QMessageBox.critical(self, 
                                obtener_traduccion('open_pptx_error_title', self.current_language),
                                obtener_traduccion('open_pptx_error_message', self.current_language).format(error=e)) 

    # --- NUEVA FUNCIÓN para manejar la cancelación ---
    def cancel_image_generation(self):
        print("Señal de cancelación recibida.")
        if hasattr(self, 'worker') and self.worker and self.worker.isRunning():
            print("Solicitando interrupción del worker...")
            self.worker.requestInterruption()
            # Opcional: podrías esperar un poco aquí, pero `deleteLater` ya está conectado a finished
            # self.worker.quit()
            # self.worker.wait(1000) # Esperar máx 1 seg
        else:
            print("Worker no encontrado o no está corriendo.")
        # El diálogo de progreso se cierra automáticamente al cancelar
        
    def update_preview_temporarily(self):
        """Actualizar la vista previa temporalmente sin guardar los cambios"""
        if not self.parent_widget or not isinstance(self.parent_widget, VentanaVistaPrevia):
            return
            
        # Marcar que se ha actualizado la vista previa temporalmente
        self.temp_preview_updated = True
        
        # Guardar los datos originales para restaurarlos si se cancela
        self.original_image_path = self.parent_widget.all_slides_data[self.parent_widget.current_slide_index].get('imagen_path')
        
        # Crear datos temporales para la vista previa
        temp_data = {
            'imagen_path': self.new_image_path,
            'titulo': self.title_edit.text(),
            'contenido': self.content_edit.toPlainText()
        }
        
        # Guardar temporalmente los nuevos datos
        self.parent_widget.all_slides_data[self.parent_widget.current_slide_index] = temp_data
        
        # Actualizar la vista previa
        self.parent_widget.mostrar_diapositiva_actual()

    # --- Función modificada para manejar la edición y guardar en PPTX ---
    def editar_diapositiva_actual(self):
        if not PPTX_AVAILABLE:
             # Opcional: Mostrar mensaje al usuario informando que la función no está disponible
             QMessageBox.warning(self, "Función no disponible", "La biblioteca 'python-pptx' es necesaria para editar el archivo PPTX directamente. Por favor, instálala.")
             return
             
        if self.current_slide_index < 0 or self.current_slide_index >= len(self.all_slides_data):
            return # No hay diapositiva válida para editar
            
        current_data = self.all_slides_data[self.current_slide_index]
        
        # Crear y mostrar el diálogo de edición - usar el idioma actual
        dialog = EditSlideDialog(current_data, self.current_language, self)
        if dialog.exec() == QDialog.Accepted:
            # Obtener los datos editados
            edited_data = dialog.get_edited_data()
            
            # 1. Actualizar los datos en la lista interna
            self.all_slides_data[self.current_slide_index] = edited_data
            
            # 2. Refrescar la vista previa para mostrar los cambios inmediatamente
            self.mostrar_diapositiva_actual()
            
            # 3. Aplicar los cambios al archivo PPTX
            # Pasar los datos originales también para poder encontrar las shapes correctas
            self._apply_changes_to_pptx(self.current_slide_index, edited_data, current_data)
            
    # --- Nueva función para aplicar cambios al archivo PPTX ---
    # Se añade 'original_data' para buscar las shapes por contenido original
    def _apply_changes_to_pptx(self, slide_index, edited_data, original_data):
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            QMessageBox.warning(self,
                                obtener_traduccion('edit_save_pptx_error_title', self.current_language),
                                obtener_traduccion('edit_save_pptx_not_found_message', self.current_language))
            return

        try:
            prs = pptx.Presentation(self.pptx_path)

            # Asegurarse que el índice es válido para la presentación
            if slide_index < 0 or slide_index >= len(prs.slides):
                 raise ValueError(f"Índice de diapositiva ({slide_index}) fuera de rango.")

            slide = prs.slides[slide_index]

            # --- Obtener la configuración de formato actualizada si está disponible ---
            format_settings = edited_data.get('format_settings')
            if format_settings:
                # Actualizar la configuración de formato en la vista previa
                self.update_format_settings(
                    format_settings['title_font_name'],
                    format_settings['content_font_name'],
                    format_settings['title_font_size'],
                    format_settings['content_font_size'],
                    format_settings['title_bold'],
                    format_settings['title_italic'],
                    format_settings['title_underline'],
                    format_settings['content_bold'],
                    format_settings['content_italic'],
                    format_settings['content_underline']
                )
                print(f"Configuración de formato actualizada: {format_settings}")
            else:
                print("No se encontró configuración de formato en los datos editados")

            # --- Actualizar Texto ---
            title_shape = None
            content_shape = None
            original_title = original_data.get('titulo', '').strip()
            original_content = original_data.get('contenido', '').strip()

            # Buscar shapes (sin cambios)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text_frame.text.strip()
                    if not title_shape and shape_text == original_title:
                        title_shape = shape
                        print(f"Shape de título encontrada por texto: '{original_title}'")
                    elif not content_shape and shape_text == original_content:
                        content_shape = shape
                        print(f"Shape de contenido encontrada por texto: '{original_content[:50]}...'")
                if title_shape and content_shape:
                    break
            
            # --- Actualizar Título (Lógica simplificada similar a la anterior, suele ser menos complejo) ---
            if title_shape:
                print(f"Actualizando título en shape: {title_shape.name if hasattr(title_shape, 'name') else 'N/A'}")
                # --- Refined Original Format Detection ---
                original_alignment = PP_ALIGN.LEFT # Default
                original_color_rgb = RGBColor(0, 0, 0) # Default BLACK now, safer than white
                title_format_detected = False
                try:
                    if title_shape and title_shape.has_text_frame and title_shape.text_frame.paragraphs:
                        first_para = title_shape.text_frame.paragraphs[0]
                        original_alignment = first_para.alignment or PP_ALIGN.LEFT

                        # Try run color first
                        if first_para.runs and hasattr(first_para.runs[0].font.color, 'type'):
                            if first_para.runs[0].font.color.type == MSO_COLOR_TYPE.RGB:
                                original_color_rgb = first_para.runs[0].font.color.rgb
                                title_format_detected = True
                                print(f"Title original - Color RGB from run: {original_color_rgb}")
                            # OPTIONAL: Handle THEME color if needed, though we can't directly reapply it easily
                            # elif first_para.runs[0].font.color.type == MSO_COLOR_TYPE.THEME:
                            #     print(f"Title original - Color from run is THEME: {first_para.runs[0].font.color.theme_color}")

                        # If run color wasn't RGB, try paragraph color
                        if not title_format_detected and hasattr(first_para.font.color, 'type'):
                            if first_para.font.color.type == MSO_COLOR_TYPE.RGB:
                                original_color_rgb = first_para.font.color.rgb
                                title_format_detected = True
                                print(f"Title original - Color RGB from paragraph: {original_color_rgb}")
                            # elif first_para.font.color.type == MSO_COLOR_TYPE.THEME:
                            #     print(f"Title original - Color from paragraph is THEME: {first_para.font.color.theme_color}")

                        if not title_format_detected:
                            print(f"Title original - Could not detect RGB color, using default BLACK.")

                        print(f"Title original detected - Alineación: {original_alignment}, Color RGB: {original_color_rgb}")
                    else:
                        print(f"Title original - Title shape or text frame invalid, using defaults.")

                except Exception as e:
                    print(f"WARN: No se pudo obtener formato original título: {e}, usando defaults.")
                    original_color_rgb = RGBColor(0, 0, 0) # Ensure default is black on error

                try:
                    nuevo_titulo = edited_data['titulo']
                    tf = title_shape.text_frame
                    tf.clear()
                    lines = nuevo_titulo.split('\n')
                    p = tf.add_paragraph()
                    p.text = lines[0] if lines else ""
                    for line in lines[1:]:
                        p = tf.add_paragraph()
                        p.text = line
                    
                    # Apply formatting (uses potentially improved original_color_rgb)
                    for para_idx, para in enumerate(tf.paragraphs):
                        para.alignment = original_alignment
                        
                        # Usar la configuración de formato del editor si está disponible
                        if format_settings:
                            # Apply font properties directly to paragraph first
                            para.font.name = format_settings['title_font_name']
                            para.font.size = Pt(format_settings['title_font_size'])
                            para.font.bold = format_settings['title_bold']
                            para.font.italic = format_settings['title_italic']
                            para.font.underline = format_settings['title_underline']
                        else:
                            # Usar la configuración actual de la vista previa
                            para.font.name = self.title_font_name
                            para.font.size = Pt(self.title_font_size)
                            para.font.bold = self.title_bold
                            para.font.italic = self.title_italic
                            para.font.underline = self.title_underline
                        
                        try:
                            # Apply the detected or default color
                            para.font.color.rgb = original_color_rgb
                            # print(f"  Title Para {para_idx}: Applied color {original_color_rgb}")
                        except Exception as color_err:
                            print(f"WARN: Could not apply title color to paragraph {para_idx}: {color_err}")
                            # Apply default black if setting detected color failed
                            try:
                                para.font.color.rgb = RGBColor(0, 0, 0)
                            except: pass # Ignore if setting black also fails

                    print(f"Título actualizado a: '{nuevo_titulo}'")
                except Exception as e:
                    print(f"Error fatal al procesar título: {e}")
                    QMessageBox.critical(self, "Error PPTX", f"Error al actualizar el título: {e}")
            else:
                print("ERROR: No se encontró la shape del título para actualizar.")

            # --- Actualizar Contenido (Lógica Mejorada) ---
            if content_shape:
                print(f"Actualizando contenido en shape: {content_shape.name if hasattr(content_shape, 'name') else 'N/A'}")
                original_alignments = []
                original_colors_rgb = []
                original_margins = { "left": None, "right": None, "top": None, "bottom": None }
                original_vertical_anchor = MSO_ANCHOR.TOP # Default
                
                # --- 1. Guardar formato original detallado --- 
                try:
                    original_tf = content_shape.text_frame
                    # Guardar márgenes y anclaje del text_frame
                    original_margins["left"] = original_tf.margin_left
                    original_margins["right"] = original_tf.margin_right
                    original_margins["top"] = original_tf.margin_top
                    original_margins["bottom"] = original_tf.margin_bottom
                    original_vertical_anchor = original_tf.vertical_anchor or MSO_ANCHOR.TOP
                    print(f"Contenido original - Márgenes: {original_margins}, Anclaje: {original_vertical_anchor}")

                    # Guardar formato de cada párrafo original
                    for p_idx, para in enumerate(original_tf.paragraphs):
                        # Guardar alineación
                        alignment = para.alignment or PP_ALIGN.LEFT # Default a Izquierda si no está definida
                        original_alignments.append(alignment)
                        
                        # --- Guardar color (REFINADO) --- 
                        color_rgb = None
                        # Intentar obtener del primer run con color RGB explícito
                        found_run_color = False
                        if para.runs:
                            for run in para.runs:
                                try:
                                    if hasattr(run.font.color, 'type') and run.font.color.type == MSO_COLOR_TYPE.RGB:
                                        color_rgb = run.font.color.rgb
                                        found_run_color = True
                                        # print(f"DEBUG: Color RGB encontrado en run: {color_rgb}")
                                        break # Usar el primer color RGB encontrado en los runs
                                except AttributeError:
                                    # Puede que el run no tenga font o color definidos
                                    continue 
                        
                        # Si no se encontró en los runs, intentar obtener del párrafo
                        if not found_run_color:
                            try:
                                if hasattr(para.font.color, 'type') and para.font.color.type == MSO_COLOR_TYPE.RGB:
                                    color_rgb = para.font.color.rgb
                                    # print(f"DEBUG: Color RGB encontrado en párrafo: {color_rgb}")
                                else:
                                    # Opcional: Comprobar si es color de tema, aunque no podamos usarlo directamente
                                    if hasattr(para.font.color, 'type') and para.font.color.type == MSO_COLOR_TYPE.THEME:
                                         print(f"WARN: Contenido Para {p_idx} - Color es de TEMA, no RGB. No se puede preservar directamente.")
                                    else:
                                         print(f"WARN: Contenido Para {p_idx} - No se pudo obtener color RGB (ni runs ni párrafo). Tipo: {getattr(para.font.color, 'type', 'N/A')}")
                            except AttributeError:
                                print(f"WARN: Contenido Para {p_idx} - Atributo de color no encontrado en párrafo.")
                        
                        original_colors_rgb.append(color_rgb) # Puede ser None
                        # print(f"Contenido Para {p_idx} original - Alineación: {alignment}, Color Guardado: {color_rgb}")
                        
                except Exception as e:
                    print(f"WARN: No se pudo obtener el formato original completo del contenido: {e}")
                    import traceback
                    traceback.print_exc()

                # --- 2. Limpiar y recrear párrafos --- 
                try:
                    nuevo_contenido = edited_data['contenido']
                    tf = content_shape.text_frame
                    tf.clear()
                    lines = nuevo_contenido.split('\n')
                    # Crear los párrafos nuevos
                    new_paragraphs = []
                    p = tf.add_paragraph()
                    p.text = lines[0] if lines else ""
                    new_paragraphs.append(p)
                    for line in lines[1:]:
                        p = tf.add_paragraph()
                        p.text = line
                        new_paragraphs.append(p)

                    # --- 3. Reaplicar formato específico --- 
                    # Reaplicar márgenes y anclaje al text_frame
                    if original_margins["left"] is not None: tf.margin_left = original_margins["left"]
                    if original_margins["right"] is not None: tf.margin_right = original_margins["right"]
                    if original_margins["top"] is not None: tf.margin_top = original_margins["top"]
                    if original_margins["bottom"] is not None: tf.margin_bottom = original_margins["bottom"]
                    tf.vertical_anchor = original_vertical_anchor
                    print(f"Márgenes y anclaje reaplicados al TextFrame.")

                    # Aplicar formato a cada párrafo nuevo
                    num_original_paras = len(original_alignments)
                    default_color = RGBColor(0, 0, 0) # Negro por defecto
                    
                    for para_idx, para in enumerate(new_paragraphs):
                        # --- Determinar Alineación y Color a Aplicar ---
                        alignment_to_apply = PP_ALIGN.LEFT # Default
                        if para_idx < num_original_paras:
                            alignment_to_apply = original_alignments[para_idx]
                        elif num_original_paras > 0:
                            alignment_to_apply = original_alignments[-1]

                        color_to_apply = default_color # Default
                        # Usar color detectado si existe y no es None
                        if para_idx < len(original_colors_rgb) and original_colors_rgb[para_idx] is not None:
                            color_to_apply = original_colors_rgb[para_idx]
                        elif original_colors_rgb and original_colors_rgb[-1] is not None:
                             color_to_apply = original_colors_rgb[-1]
                             if para_idx >= len(original_colors_rgb): # Informar si se usa fallback del último párrafo
                                 print(f"WARN: Contenido Para {para_idx} - Usando color del último párrafo original como fallback.")
                        elif color_to_apply == default_color:
                             print(f"WARN: Contenido Para {para_idx} - Usando color NEGRO de fallback (no se detectó RGB original).")

                        # --- Aplicar Formato a Nivel de Párrafo --- 
                        para.alignment = alignment_to_apply
                        
                        # Usar la configuración de formato del editor si está disponible
                        if format_settings:
                            para.font.name = format_settings['content_font_name']
                            para.font.size = Pt(format_settings['content_font_size'])
                            para.font.bold = format_settings['content_bold']
                            para.font.italic = format_settings['content_italic']
                            para.font.underline = format_settings['content_underline']
                        else:
                            # Usar la configuración actual de la vista previa
                            para.font.name = self.content_font_name
                            para.font.size = Pt(self.content_font_size)
                            para.font.bold = self.content_bold
                            para.font.italic = self.content_italic
                            para.font.underline = self.content_underline
                        
                        para.font.color.rgb = color_to_apply

                    # Optional: Adjust text frame shape to fit text?
                    # Considerar si MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT es deseable aquí
                    # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT 
                    print(f"Contenido actualizado y formato reaplicado para: '{nuevo_contenido[:50]}...'")

                except Exception as e:
                    print(f"Error fatal al procesar contenido: {e}")
                    import traceback
                    traceback.print_exc() # Imprimir stack trace para depuración
                    QMessageBox.critical(self, "Error PPTX", f"Error al actualizar el contenido: {e}")
            else:
                print("ERROR: No se encontró la shape del contenido para actualizar.")


            # --- ACTUALIZADO: Manejar imágenes de fondo en diseños especiales --- 
            new_image_path = edited_data['imagen_path']
            original_image_path = original_data.get('imagen_path')
            
            # Verificar si es un diseño con imagen de fondo (identificado por shape name)
            background_image_shape = None
            for shp in slide.shapes:
                # Buscar la imagen identificada como fondo
                if hasattr(shp, 'name') and shp.name == "SLIDE_BACKGROUND_IMAGE":
                    background_image_shape = shp
                    print(f"Imagen de fondo encontrada: {shp.name}")
                    break
                    
            if background_image_shape:
                # Es un diseño especial con imagen de fondo (design3, design4 o design7)
                print(f"Detectado diseño especial con imagen de fondo")

                # --- Definir función auxiliar para detectar texto blanco ---
                def text_is_white(shape_obj_local):
                    if not shape_obj_local or not shape_obj_local.has_text_frame:
                        return False
                    text_frame = shape_obj_local.text_frame
                    for para in text_frame.paragraphs:
                        # Verificar color a nivel de párrafo
                        try:
                            if hasattr(para.font.color, 'type') and para.font.color.type == MSO_COLOR_TYPE.RGB:
                                if para.font.color.rgb == RGBColor(255, 255, 255):
                                    return True # Encontrado a nivel de párrafo
                        except AttributeError:
                            pass # Seguir a runs

                        # Verificar color a nivel de runs
                        if para.runs:
                            for run_idx, run in enumerate(para.runs): # Iterar todos los runs
                                try:
                                    if hasattr(run.font.color, 'type') and run.font.color.type == MSO_COLOR_TYPE.RGB:
                                        if run.font.color.rgb == RGBColor(255, 255, 255):
                                            return True # Encontrado a nivel de run
                                except AttributeError:
                                    pass # Siguiente run
                    return False # No se encontró texto blanco

                # --- Determinar is_design7 de forma más robusta ---
                is_design7 = False
                if title_shape and content_shape: # Asegurarse que las shapes de título y contenido existen
                    if text_is_white(title_shape) and text_is_white(content_shape):
                        is_design7 = True
                        print("INFO: 'is_design7' establecido a True (texto blanco en título y contenido).")
                
                # Solo reemplazar si la imagen cambió y existe
                if os.path.exists(new_image_path) and new_image_path != original_image_path:
                    print(f"La ruta de la imagen de fondo cambió a: {new_image_path}. Reemplazando...")
                    
                    # Obtener posición y tamaño de la imagen original
                    left_img_bg, top_img_bg = Inches(0), Inches(0)  # Siempre posición 0,0 para fondos
                    # Obtener las dimensiones de la presentación actual
                    width_img_bg = prs.slide_width
                    height_img_bg = prs.slide_height
                    
                    # Para design7, procesar oscurecimiento de la imagen antes de reemplazar
                    if slide_index >= 0 and slide_index < len(self.all_slides_data):
                        try:
                            if is_design7: # Usar la variable is_design7 determinada por la nueva lógica
                                print("Diseño 7 detectado (para oscurecer imagen): procesando imagen oscurecida...")
                                # Mismo procesamiento que en design7
                                APP_DATA_DIR_IMG = os.path.join(os.getenv('APPDATA'), 'Powerpoineador')
                                IMAGES_DIR_IMG = os.path.join(APP_DATA_DIR_IMG, 'images')
                                if not os.path.exists(IMAGES_DIR_IMG):
                                    os.makedirs(IMAGES_DIR_IMG, exist_ok=True)
                                from PIL import Image, ImageEnhance # Asegurar importación local si es necesario
                                image_pil = Image.open(new_image_path)
                                enhancer = ImageEnhance.Brightness(image_pil)
                                image_darker = enhancer.enhance(0.5)
                                darker_path = os.path.join(IMAGES_DIR_IMG, 'Slide_darker_edited.jpg') # Nombre diferente para evitar conflictos
                                image_darker.save(darker_path)
                                # Usar la imagen oscurecida en su lugar
                                new_image_path = darker_path
                        except Exception as e_img_proc:
                            print(f"Error al procesar imagen oscurecida para design7: {e_img_proc}")
                    
                    # ---- MÉTODO MEJORADO: Crear una nueva diapositiva con el mismo diseño y contenido ----
                    try:
                        # 1. Primero, crear una copia de seguridad del contenido de texto
                        # title_text = "" # No se usa directamente
                        # content_text = "" # No se usa directamente
                        # title_shape y content_shape ya están definidos y buscados antes
                        
                        # 2. Verificar qué diseño estamos usando por marcadores específicos
                        design_number = 3 # Por defecto, si no es design7 y no es claramente design4, asumimos design3
                        
                        if is_design7: # Usar la variable is_design7 robusta
                            design_number = 7
                        elif title_shape and content_shape:
                            # design4 tiene el texto claramente en la mitad izquierda (p.ej., primer 45%)
                            # design3 tiene el texto más cerca del centro o en la mitad derecha.
                            if title_shape.left < (prs.slide_width * 0.45): # Umbral para design4
                                design_number = 4
                            # else: design_number remains 3 (cubriendo el caso de design3)
                        
                        print(f"INFO: Recreando diapositiva con design_number: {design_number}")
                        
                        # 3. Importar el módulo de diseños
                        from Diseños_diapositivas import Diapositivas
                        
                        # 4. MÉTODO CORREGIDO: En lugar de eliminar directamente la diapositiva,
                        # vamos a crear una nueva y luego sobrescribir los contenidos
                        # El enfoque anterior intentaba manipular directamente la estructura XML
                        # pero la API no nos da acceso de esa manera
                        
                        # Crear una instancia de Diapositivas con la misma configuración que la actual
                        diapositiva = Diapositivas(prs, 
                                                  title_font_name=self.title_font_name, 
                                                  content_font_name=self.content_font_name,
                                                  title_font_size=self.title_font_size, 
                                                  content_font_size=self.content_font_size,
                                                  title_bold=self.title_bold, 
                                                  title_italic=self.title_italic, 
                                                  title_underline=self.title_underline,
                                                  content_bold=self.content_bold, 
                                                  content_italic=self.content_italic, 
                                                  content_underline=self.content_underline)
                        
                        # 5. Crear una nueva diapositiva temporal con el diseño y contenido actualizados
                        # Añadir la nueva diapositiva al final
                        if design_number == 3:
                            diapositiva.design3(None, edited_data['titulo'], edited_data['contenido'], new_image_path)
                        elif design_number == 4:
                            diapositiva.design4(None, edited_data['titulo'], edited_data['contenido'], new_image_path)
                        elif design_number == 7:
                            diapositiva.design7(None, edited_data['titulo'], edited_data['contenido'], new_image_path)
                            
                        # 6. La nueva diapositiva ahora está al final, obtenemos su índice
                        new_slide_idx = len(prs.slides) - 1
                        new_slide = prs.slides[new_slide_idx]
                        
                        # 7. Usar el método oficial de python-pptx para mover la diapositiva
                        # Mover la nueva diapositiva a la posición de la antigua
                        xml_slides = prs.slides._sldIdLst  # Este es el contenedor XML de slide IDs
                        
                        # Si la posición de destino no es el final, debemos mover la diapositiva
                        if slide_index < new_slide_idx:
                            # Obtener el elemento XML de la nueva diapositiva
                            sldId = xml_slides[new_slide_idx]  # El elemento XML que referencia la diapositiva
                            
                            # Eliminar la referencia desde su posición actual
                            xml_slides.remove(sldId)
                            
                            # Insertar en la posición deseada
                            xml_slides.insert(slide_index, sldId)
                        
                        # 8. Eliminar la diapositiva original (ahora desplazada al siguiente índice)
                        # Ya que agregamos una nueva al final y posiblemente la movimos,
                        # la antigua estará en slide_index+1 si movimos la nueva, o en slide_index si no
                        # Usamos el enfoque oficial de python-pptx para eliminar la diapositiva
                        idx_to_remove = slide_index
                        if slide_index < new_slide_idx:  # Si movimos la nueva a la posición original
                            idx_to_remove = slide_index + 1  # La antigua está una posición más allá
                            
                        if idx_to_remove < len(prs.slides):
                            # Obtener el ID de la diapositiva a eliminar
                            try:
                                sldId_to_remove = xml_slides[idx_to_remove]
                                xml_slides.remove(sldId_to_remove)
                                print(f"Diapositiva antigua eliminada en la posición {idx_to_remove+1}")
                            except Exception as e_remove:
                                print(f"Error al eliminar diapositiva antigua: {e_remove}")
                                # No es crítico si falla la eliminación, simplemente tendríamos una diapositiva duplicada
                                
                        print(f"Nueva diapositiva creada con éxito en la posición {slide_index+1}")
                        
                    except Exception as e:
                        print(f"ERROR al recrear diapositiva: {e}")
                        import traceback
                        traceback.print_exc()
                        QMessageBox.warning(self, "Error al reemplazar imagen", 
                                          f"Se produjo un error al reemplazar la imagen de fondo: {e}")
                    
                elif not os.path.exists(new_image_path):
                    print(f"Error: La nueva imagen no existe en la ruta: {new_image_path}")
                    QMessageBox.warning(self, "Error Imagen", f"La ruta de la nueva imagen no es válida:\n{new_image_path}")
                else:
                    print("La imagen de fondo no ha cambiado. No se reemplaza.")
            else:
                # --- Caso de imágenes regulares (método anterior) ---
                pic_shape = None
                for shp in slide.shapes:
                    # Asegurarse de que no sea un placeholder de imagen vacío
                    if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        pic_shape = shp
                        print(f"Shape de imagen encontrada: {pic_shape.name if hasattr(pic_shape, 'name') else 'N/A'}")
                        break # Usar la primera que se encuentre
                    elif shp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shp, 'placeholder_format') and shp.placeholder_format.type == MSO_SHAPE_TYPE.PICTURE:
                         # Podría ser un placeholder de imagen, verificar si tiene imagen
                         if hasattr(shp, 'image'):
                             pic_shape = shp
                             print(f"Shape de imagen (placeholder) encontrada: {pic_shape.name if hasattr(pic_shape, 'name') else 'N/A'}")
                             break

                if pic_shape:
                    # ---- SOLO REEMPLAZAR SI LA IMAGEN CAMBIÓ Y ES VÁLIDA ----
                    if os.path.exists(new_image_path) and new_image_path != original_image_path:
                        print(f"La ruta de la imagen cambió a: {new_image_path}. Reemplazando...")
                        # Obtener posición y tamaño de la imagen original ANTES de borrarla
                        left, top, width, height = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height

                        # Eliminar la shape de imagen antigua
                        # Comprobar si es un placeholder antes de intentar eliminarlo de la forma estándar
                        if pic_shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
                            try:
                                # ENFOQUE MENOS INVASIVO - MANTENER COMPATIBILIDAD
                                # En lugar de manipular el XML directamente, simplemente añadimos la nueva imagen
                                # y no eliminamos la antigua, confiando en que PowerPoint la manejará correctamente
                                # cuando se abra (se mostrará la más reciente)
                                print("Añadiendo nueva imagen sobre la anterior")
                                new_pic = slide.shapes.add_picture(new_image_path, left, top, width, height)
                                
                                # Intentar renombrar la nueva imagen con el mismo nombre que la antigua si tenía alguno
                                if hasattr(pic_shape, 'name') and pic_shape.name:
                                    try:
                                        new_pic.name = pic_shape.name
                                    except:
                                        pass
                            except Exception as e_rem:
                                print(f"WARN: No se pudo gestionar la imagen: {e_rem}")
                        else:
                             print("INFO: La imagen original es un placeholder, se intentará actualizar")
                             # Para placeholders, intentamos una técnica más segura
                             try:
                                 new_pic = slide.shapes.add_picture(new_image_path, left, top, width, height)
                                 if hasattr(pic_shape, 'name') and pic_shape.name:
                                     try:
                                         new_pic.name = pic_shape.name
                                     except:
                                         pass
                             except Exception as e_add:
                                 print(f"ERROR: No se pudo añadir la nueva imagen sobre el placeholder: {e_add}")
                                 QMessageBox.warning(self, "Error Imagen", f"No se pudo reemplazar la imagen: {e_add}")
                    elif not os.path.exists(new_image_path):
                        print(f"Error: La nueva imagen no existe en la ruta: {new_image_path}")
                        QMessageBox.warning(self, "Error Imagen", f"La ruta de la nueva imagen no es válida:\n{new_image_path}")
                    else:
                        # La ruta es la misma o la nueva no existe pero es igual a la original
                        print("La imagen no ha cambiado o la nueva ruta no es válida y es igual a la original. No se reemplaza.")
                    # ---- FIN DEL BLOQUE DE REEMPLAZO ----
                else:
                    print("WARN: No se encontró ninguna shape de imagen para actualizar.")
            # --- Fin de gestión de imágenes ---

            # --- Guardar la presentación ---
            prs.save(self.pptx_path)
            print("Presentación guardada con éxito.")

        except Exception as e:
            print(f"Error general al guardar cambios en PPTX: {e}")
            import traceback
            traceback.print_exc() # Imprimir stack trace completo
            QMessageBox.critical(self,
                               obtener_traduccion('edit_save_pptx_error_title', self.current_language),
                               obtener_traduccion('edit_save_pptx_error_message', self.current_language).format(error=e))

    # --- NUEVA FUNCIÓN: Abrir presentación en el programa predeterminado del sistema ---
    def abrir_presentacion(self):
        if not self.pptx_path or not os.path.exists(self.pptx_path):
            QMessageBox.warning(self,
                                obtener_traduccion('open_pptx_error_title', self.current_language),
                                obtener_traduccion('open_pptx_not_found_message', self.current_language))
            return
            
        try:
            # En Windows usar 'start', en otros sistemas puede variar
            if sys.platform == 'win32':
                os.startfile(self.pptx_path)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', self.pptx_path])
            else:  # Linux y otros
                subprocess.call(['xdg-open', self.pptx_path])
                
            print(f"Presentación abierta: {self.pptx_path}")
        except Exception as e:
            print(f"Error al abrir la presentación: {e}")
            QMessageBox.critical(self, 
                                obtener_traduccion('open_pptx_error_title', self.current_language),
                                obtener_traduccion('open_pptx_error_message', self.current_language).format(error=e)) 

    # --- NUEVA FUNCIÓN para manejar la cancelación ---
    def cancel_image_generation(self):
        print("Señal de cancelación recibida.")
        if hasattr(self, 'worker') and self.worker and self.worker.isRunning():
            print("Solicitando interrupción del worker...")
            self.worker.requestInterruption()
            # Opcional: podrías esperar un poco aquí, pero `deleteLater` ya está conectado a finished
            # self.worker.quit()
            # self.worker.wait(1000) # Esperar máx 1 seg
        else:
            print("Worker no encontrado o no está corriendo.")
        # El diálogo de progreso se cierra automáticamente al cancelar
        
    def update_preview_temporarily(self):
        """Actualizar la vista previa temporalmente sin guardar los cambios"""
        if not self.parent_widget or not isinstance(self.parent_widget, VentanaVistaPrevia):
            return
            
        # Marcar que se ha actualizado la vista previa temporalmente
        self.temp_preview_updated = True
        
        # Guardar los datos originales para restaurarlos si se cancela
        self.original_image_path = self.parent_widget.all_slides_data[self.parent_widget.current_slide_index].get('imagen_path')
        
        # Crear datos temporales para la vista previa
        temp_data = {
            'imagen_path': self.new_image_path,
            'titulo': self.title_edit.text(),
            'contenido': self.content_edit.toPlainText()
        }
        
        # Guardar temporalmente los nuevos datos
        self.parent_widget.all_slides_data[self.parent_widget.current_slide_index] = temp_data
        
        # Actualizar la vista previa
        self.parent_widget.mostrar_diapositiva_actual()

    # Función para actualizar el mensaje cuando no hay diapositivas
    def actualizar_mensaje_vacio(self):
        try:
            mensaje = obtener_traduccion('vista_previa_empty_message', self.current_language)
            if not mensaje:
                mensaje = "Aquí aparecerán las diapositivas generadas"
            
            # Actualizar SOLO el objeto mensaje_vacio principal
            if hasattr(self, 'mensaje_vacio') and self.mensaje_vacio is not None:
                try:
                    self.mensaje_vacio.setText(mensaje)
                except RuntimeError:
                    # Ignorar si el objeto ya ha sido eliminado
                    pass
                
            # No intentar buscar en la jerarquía de widgets para evitar
            # acceder a widgets que puedan haber sido eliminados
        except Exception as e:
            print(f"Error en actualizar_mensaje_vacio: {str(e)}")

    # Función para agregar una nueva diapositiva a la vista previa  
    def agregar_diapositiva(self, imagen_path, titulo, contenido):
        # Guardar datos de la diapositiva
        if os.path.exists(imagen_path):
            # Crear configuración de formato específica para esta diapositiva
            # utilizando la configuración global actual
            format_settings = {
                'title_font_name': self.title_font_name,
                'content_font_name': self.content_font_name,
                'title_font_size': self.title_font_size,
                'content_font_size': self.content_font_size,
                'title_bold': self.title_bold,
                'title_italic': self.title_italic,
                'title_underline': self.title_underline,
                'content_bold': self.content_bold,
                'content_italic': self.content_italic,
                'content_underline': self.content_underline
            }
            
            self.all_slides_data.append({
                'imagen_path': imagen_path,
                'titulo': titulo,
                'contenido': contenido,
                'format_settings': format_settings  # Incluir configuración de formato específica
            })
            
            # Solo actualizar a la última si no estamos en modo navegación
            if not self.modo_navegacion:
                self.current_slide_index = len(self.all_slides_data) - 1
                self.mostrar_diapositiva_actual()
            else:
                # Si estamos en modo navegación, solo actualizar el contador
                # con el total actualizado pero manteniendo la diapositiva actual
                self.slide_counter_label.setText(f"{self.current_slide_index + 1}/{len(self.all_slides_data)}")
            
            # Actualizar botones de navegación
            self.actualizar_botones_navegacion() 