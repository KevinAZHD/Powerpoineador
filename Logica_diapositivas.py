import sys, os, random
from pptx import Presentation
from Diseños_diapositivas import Diapositivas
from modelos.IA_sdxl import generar_imagen as generar_imagen_sdxl
from modelos.IA_fluxschnell import generar_imagen as generar_imagen_fluxschnell
from modelos.IA_flux8 import generar_imagen as generar_imagen_flux8
from modelos.IA_flux16 import generar_imagen as generar_imagen_flux16
from modelos.IA_dgmtnzflux import generar_imagen as generar_imagen_diego
from modelos.IA_fluxpulid import generar_imagen as generar_imagen_flux
from modelos.IA_photomaker import generar_imagen as generar_imagen_photomaker
from modelos.IA_imagen3 import generar_imagen as generar_imagen_imagen3
from modelos.IA_imagen3fast import generar_imagen as generar_imagen_imagen3fast
from modelos.IA_sana import generar_imagen as generar_imagen_sana
from modelos.IA_sana_sprint import generar_imagen as generar_imagen_sana_sprint
from modelos.IA_model3_4 import generar_imagen as generar_imagen_model3_4
from modelos.IA_grok2_image import generar_imagen as generar_imagen_grok
from modelos.IA_gemini2_flash_image import generar_imagen as generar_imagen_gemini_flash
from Traducciones import obtener_traduccion

# Definir la ruta de la carpeta de datos de la aplicación según el sistema operativo
if sys.platform == 'win32':
    APP_DATA_DIR = os.path.join(os.getenv('APPDATA'), 'Powerpoineador')
elif sys.platform == 'darwin':
    APP_DATA_DIR = os.path.join(os.path.expanduser('~'), 'Library', 'Application Support', 'Powerpoineador')
else:
    APP_DATA_DIR = os.path.join(os.path.expanduser('~'), '.Powerpoineador')

# Crear carpeta de imágenes si no existe
IMAGES_DIR = os.path.join(APP_DATA_DIR, 'images')
if not os.path.exists(IMAGES_DIR):
    os.makedirs(IMAGES_DIR)

# Función para obtener respuesta del modelo con reintentos
def obtener_respuesta_ia(descripcion, modelo, signals=None):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir un mensaje indicando que se está intentando generar una respuesta
        log_message(obtener_traduccion('intentando_generar_respuesta', current_language).format(modelo=modelo))
        
        # Verificar cuál es el modelo que se está utilizando
        if modelo == 'meta-llama-3.1-405b-instruct [$0.0067]':
            from modelos.IA_llama3 import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'meta-llama-4-scout-instruct [$0.00046]':
            from modelos.IA_llama4s import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'meta-llama-4-maverick-instruct [$0.00067]':
            from modelos.IA_llama4m import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'claude-3.7-sonnet [$0.0105]':
            from modelos.IA_sonnet3_7 import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'claude-3.5-sonnet [$0.01312]':
            from modelos.IA_sonnet3_5 import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'claude-3.5-haiku [$0.0035]':
            from modelos.IA_haiku import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'grok-2-1212':
            from modelos.IA_grok2 import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'grok-3-beta':
            from modelos.IA_grok3 import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'grok-3-mini-beta':
            from modelos.IA_grok3_mini import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'gemini-2.5-pro-exp-03-25':
            from modelos.IA_gemini2_pro import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'gemini-2.0-flash-thinking-exp-01-21':
            from modelos.IA_gemini2_flash_thinking import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        elif modelo == 'deepseek-r1 [$0.007]':
            from modelos.IA_deepseek import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)
        else:
            from modelos.IA_dolphin import intentar_obtener_respuesta
            respuesta = intentar_obtener_respuesta(descripcion, signals)

        # Verificar si se pudo obtener respuesta del modelo
        if respuesta:
            log_message(obtener_traduccion('respuesta_modelo', current_language).format(modelo=modelo, respuesta=respuesta))
            log_message(obtener_traduccion('tupla_generada', current_language).format(respuesta=respuesta))
            return respuesta
        
        return None
        
    except Exception as e:
        # Imprimir un mensaje indicando que ocurrió un error al obtener la respuesta
        log_message(obtener_traduccion('error_obtener_respuesta', current_language).format(error=str(e)))
        return None
    finally:
        import gc
        gc.collect()

# Función para generar texto (título y contenido) con un modelo de IA
def generar_texto_ia(tipo, contenido_actual, descripcion, modelo, signals=None):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Crear un prompt específico según si estamos generando título o contenido
        if tipo == "titulo":
            prompt = obtener_traduccion('prompt_titulo', current_language).format(
                descripcion=descripcion, 
                contenido_actual=contenido_actual
            )
            log_message(obtener_traduccion('generando_titulo', current_language).format(modelo=modelo))
        else:  # contenido
            prompt = obtener_traduccion('prompt_contenido', current_language).format(
                descripcion=descripcion, 
                contenido_actual=contenido_actual
            )
            log_message(obtener_traduccion('generando_contenido', current_language).format(modelo=modelo))
        
        # Verificar cuál es el modelo que se está utilizando
        if modelo == 'meta-llama-3.1-405b-instruct [$0.0067]':
            from modelos.IA_llama3 import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'meta-llama-4-scout-instruct [$0.00046]':
            from modelos.IA_llama4s import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'meta-llama-4-maverick-instruct [$0.00067]':
            from modelos.IA_llama4m import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'claude-3.7-sonnet [$0.0105]':
            from modelos.IA_sonnet3_7 import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'claude-3.5-sonnet [$0.01312]':
            from modelos.IA_sonnet3_5 import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'claude-3.5-haiku [$0.0035]':
            from modelos.IA_haiku import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'grok-2-1212':
            from modelos.IA_grok2 import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'grok-3-beta':
            from modelos.IA_grok3 import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'grok-3-mini-beta':
            from modelos.IA_grok3_mini import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'gemini-2.5-pro-exp-03-25':
            from modelos.IA_gemini2_pro import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'gemini-2.0-flash-thinking-exp-01-21':
            from modelos.IA_gemini2_flash_thinking import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        elif modelo == 'deepseek-r1 [$0.007]':
            from modelos.IA_deepseek import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)
        else:
            from modelos.IA_dolphin import generar_texto_simple
            texto_generado = generar_texto_simple(prompt, signals)

        # Verificar si se pudo obtener respuesta del modelo
        if texto_generado:
            # Imprimir un mensaje con el texto generado
            if tipo == "titulo":
                log_message(obtener_traduccion('titulo_generado', current_language).format(texto=texto_generado))
            else:
                log_message(obtener_traduccion('contenido_generado', current_language).format(texto=texto_generado))
            return texto_generado
        
        return None
        
    except Exception as e:
        # Imprimir un mensaje indicando que ocurrió un error
        if tipo == "titulo":
            log_message(obtener_traduccion('error_generar_titulo', current_language).format(error=str(e)))
        else:
            log_message(obtener_traduccion('error_generar_contenido', current_language).format(error=str(e)))
        return None
    finally:
        import gc
        gc.collect()

# Función para generar una imagen con un modelo de IA
def generar_imagen_ia(section, content, descripcion, modelo, signals=None):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Ya no imprimir mensaje aquí, ya se hace en generar_presentacion
        
        if modelo == 'sdxl-lightning-4step [$0.0014]':
            return generar_imagen_sdxl(section, content, descripcion, signals, False)
        elif modelo == 'flux-schnell [$0.003]':
            return generar_imagen_fluxschnell(section, content, descripcion, signals, False)
        elif modelo == 'hyper-flux-8step [$0.027]':
            return generar_imagen_flux8(section, content, descripcion, signals, False)
        elif modelo == 'photomaker [$0.0070]':
            return generar_imagen_photomaker(section, content, descripcion, None, signals, False)
        elif modelo == 'hyper-flux-16step [$0.020]':
            return generar_imagen_flux16(section, content, descripcion, signals, False)
        elif modelo == 'dgmtnzflux [$0.03]':
            return generar_imagen_diego(section, content, descripcion, signals, False)
        elif modelo == 'sana [$0.0015]':
            return generar_imagen_sana(section, content, descripcion, signals, False)
        elif modelo == 'sana-sprint-1.6b [$0.0015]':
            return generar_imagen_sana_sprint(section, content, descripcion, signals, False)
        elif modelo == 'imagen-3 [$0.05]':
            return generar_imagen_imagen3(section, content, descripcion, signals, False)
        elif modelo == 'imagen-3-fast [$0.025]':
            return generar_imagen_imagen3fast(section, content, descripcion, signals, False)
        elif modelo == 'model3_4 [$0.00098]':
            return generar_imagen_model3_4(section, content, descripcion, signals, False)
        elif modelo == 'grok-2-image-1212':
            return generar_imagen_grok(section, content, descripcion, signals, False)
        elif modelo == 'gemini-2.0-flash-exp-image-generation':
            return generar_imagen_gemini_flash(section, content, descripcion, signals, False)
        else:
            return generar_imagen_flux(section, content, descripcion, None, signals, False)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una presentación con un modelo de IA
def generar_presentacion(modelo_texto, modelo_imagen, descripcion, auto_open, imagen_personalizada, filename, signals=None, title_font_name='Calibri', content_font_name='Calibri', title_font_size=16, content_font_size=10, title_bold=False, title_italic=False, title_underline=False, content_bold=False, content_italic=False, content_underline=False, disenos_aleatorios=True, selected_layout_index=1):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Crear una nueva presentación
        presentation = Presentation()
        # Crear un objeto para aplicar diseños a las diapositivas, pasando la fuente y tamaños
        slide_designs = Diapositivas(presentation, title_font_name, content_font_name, title_font_size, content_font_size, title_bold, title_italic, title_underline, content_bold, content_italic, content_underline)

        # Imprimir un mensaje indicando que se está generando texto con el modelo
        log_message(obtener_traduccion('generando_texto', current_language).format(modelo=modelo_texto))
        
        # Obtener la respuesta del modelo
        respuesta = obtener_respuesta_ia(descripcion, modelo_texto, signals)
        
        # Verificar si se pudo obtener respuesta del modelo
        if not respuesta:
            raise Exception(obtener_traduccion('no_respuesta_modelo_texto', current_language))
        
        # Obtener las secciones del contenido
        sections = respuesta
        # Inicializar el número de diapositiva
        slide_number = 1
        # Obtener el número total de diapositivas
        total_slides = len(sections)
        # Crear una lista para almacenar las imágenes generadas
        imagenes_generadas = []
        
        # Iterar sobre las secciones del contenido
        for section, content in sections.items():
            # Imprimir un mensaje indicando que se está generando una imagen
            log_message(obtener_traduccion('generando_imagen', current_language).format(numero=slide_number, total=total_slides))
            try:
                # Verificar cuál es el modelo que se está utilizando
                if modelo_imagen in ['flux-pulid [$0.027]', 'photomaker [$0.0011]']:
                    if imagen_personalizada and os.path.exists(imagen_personalizada):
                        if modelo_imagen == 'flux-pulid [$0.027]':
                            img = generar_imagen_flux(section, content, descripcion, imagen_personalizada, signals, False)
                        else:  # photomaker
                            img = generar_imagen_photomaker(section, content, descripcion, imagen_personalizada, signals, False)
                    else:
                        raise RuntimeError(obtener_traduccion('error_imagen_requerida', current_language))
                else:
                    img = generar_imagen_ia(section, content, descripcion, modelo_imagen, signals)
                
                # Guardar la imagen generada en la carpeta de imágenes
                imagen_path = os.path.join(IMAGES_DIR, f"Slide{slide_number}.jpg")
                img.save(imagen_path)
                # Agregar la ruta de la imagen generada a la lista de imágenes generadas
                imagenes_generadas.append(imagen_path)
                # Imprimir un mensaje indicando que la imagen se generó correctamente
                log_message(obtener_traduccion('imagen_generada_correctamente', current_language).format(numero=slide_number))
                
                # Emitir señal para actualizar la vista previa con imagen y texto
                if signals:
                    signals.nueva_diapositiva.emit(imagen_path, section, content)
                
                if signals:
                    # Emitir una señal para actualizar el progreso
                    signals.update_progress.emit(slide_number, total_slides)
                
                slide_number += 1
            except Exception as e:
                # Imprimir un mensaje indicando que ocurrió un error al generar la imagen
                log_message(obtener_traduccion('error_generando_imagen', current_language).format(numero=slide_number, error=str(e)))
                raise

        # Imprimir un mensaje indicando que se está aplicando diseños a las diapositivas
        log_message(obtener_traduccion('aplicando_disenos', current_language))
        # Crear una lista con los diseños disponibles
        designs = list(range(1, 10))
        # Modificar la lógica de diseños
        if selected_layout_index == 0:  # Aleatorio
            random.shuffle(designs)
            disenos_aleatorios = True
        else:
            disenos_aleatorios = False
        # Iterar sobre las secciones del contenido
        for i, (section, content) in enumerate(sections.items()):
            try:
                # Verificar si es la primera diapositiva
                if i == 0:
                    # Aplicar diseño de introducción (design0)
                    slide_designs.design0(presentation.slides, section, content, imagenes_generadas[i])
                # LÓGICA DE APLICACIÓN DE DISEÑO MODIFICADA
                elif disenos_aleatorios:
                    # Si la aleatorización está activada, usar la lista (barajada o no)
                    design_index = (i - 1) % len(designs)
                    design = designs[design_index]
                    # Aplicar diseño seleccionado aleatoriamente
                    if design == 1:
                        slide_designs.design1(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 2:
                        slide_designs.design2(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 3:
                        slide_designs.design3(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 4:
                        slide_designs.design4(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 5:
                        slide_designs.design5(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 6:
                        slide_designs.design6(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 7:
                        slide_designs.design7(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 8:
                        slide_designs.design8(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 9:
                        slide_designs.design9(presentation.slides, section, content, imagenes_generadas[i])
                elif selected_layout_index == 5: # Minimalista (fijo)
                    slide_designs.design9(presentation.slides, section, content, imagenes_generadas[i])
                elif selected_layout_index == 1: # Formal (fijo)
                    slide_designs.design8(presentation.slides, section, content, imagenes_generadas[i])
                elif selected_layout_index == 8: # Visual (fijo)
                    # Alternar entre design1 y design2 para crear un diseño visual dinámico
                    if i % 2 == 0:
                        slide_designs.design1(presentation.slides, section, content, imagenes_generadas[i])
                    else:
                        slide_designs.design2(presentation.slides, section, content, imagenes_generadas[i])
                elif selected_layout_index == 3: # Comparacion (fijo)
                    # Alternar entre design5 y design6 para crear un estilo de comparación
                    if i % 2 == 0:
                        slide_designs.design5(presentation.slides, section, content, imagenes_generadas[i])
                    else:
                        slide_designs.design6(presentation.slides, section, content, imagenes_generadas[i])
                else: # Si no es aleatorio y no es 1, 5, 8 ni 3, usar el índice seleccionado (que debería ser 6 u otro futuro)
                      # Como fallback o si se selecciona explícitamente "Libre", aplicamos la secuencia *fija*
                    design_index = (i - 1) % len(designs) # Usar la lista *ordenada*
                    design = designs[design_index]
                    if design == 1:
                        slide_designs.design1(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 2:
                        slide_designs.design2(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 3:
                        slide_designs.design3(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 4:
                        slide_designs.design4(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 5:
                        slide_designs.design5(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 6:
                        slide_designs.design6(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 7:
                        slide_designs.design7(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 8:
                        slide_designs.design8(presentation.slides, section, content, imagenes_generadas[i])
                    elif design == 9:
                        slide_designs.design9(presentation.slides, section, content, imagenes_generadas[i])
                # Imprimir un mensaje indicando que la diapositiva se completó correctamente
                log_message(obtener_traduccion('diapositiva_completada', current_language).format(numero=i+1, total=total_slides))
            except Exception as e:
                # Imprimir un mensaje indicando que ocurrió un error al aplicar el diseño a la diapositiva
                log_message(obtener_traduccion('error_aplicando_diseno', current_language).format(numero=i+1, error=str(e)))
                raise

        # Guardar la presentación en un archivo
        presentation.save(filename)

        # Liberar memoria
        presentation = None
        slide_designs = None

        # Emitir una señal para indicar que el proceso terminó
        if signals:
            signals.finished.emit()

    except Exception as e:
        # Imprimir un mensaje indicando que ocurrió un error durante la generación de la presentación
        log_message(obtener_traduccion('error_generacion_presentacion', current_language).format(error=str(e)))
        # Emitir señal de error
        if signals:
            signals.error.emit(str(e))
    finally:
        # Liberar memoria
        import gc
        gc.collect()

# Función para generar una imagen con el modelo Flux
def generar_imagen_flux(section, content, descripcion, image1, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Flux)")
        
        # Importar la función original
        from modelos.IA_fluxpulid import generar_imagen
        return generar_imagen(section, content, descripcion, image1)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Photomaker
def generar_imagen_photomaker(section, content, descripcion, image1, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Photomaker)")
        
        # Importar la función original
        from modelos.IA_photomaker import generar_imagen
        return generar_imagen(section, content, descripcion, image1)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo SDXL
def generar_imagen_sdxl(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (SDXL)")
        
        # Importar la función original
        from modelos.IA_sdxl import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo FluxSchnell
def generar_imagen_fluxschnell(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (FluxSchnell)")
        
        # Importar la función original
        from modelos.IA_fluxschnell import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Flux8
def generar_imagen_flux8(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Flux8)")
        
        # Importar la función original
        from modelos.IA_flux8 import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Flux16
def generar_imagen_flux16(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Flux16)")
        
        # Importar la función original
        from modelos.IA_flux16 import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Diego
def generar_imagen_diego(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Diego)")
        
        # Importar la función original
        from modelos.IA_dgmtnzflux import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Sana
def generar_imagen_sana(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Sana)")
        
        # Importar la función original
        from modelos.IA_sana import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo SanaSprint
def generar_imagen_sana_sprint(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (SanaSprint)")
        
        # Importar la función original
        from modelos.IA_sana_sprint import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Imagen3
def generar_imagen_imagen3(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Imagen3)")
        
        # Importar la función original
        from modelos.IA_imagen3 import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Imagen3Fast
def generar_imagen_imagen3fast(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Imagen3Fast)")
        
        # Importar la función original
        from modelos.IA_imagen3fast import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Model3_4
def generar_imagen_model3_4(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Model3_4)")
        
        # Importar la función original
        from modelos.IA_model3_4 import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))

# Función para generar una imagen con el modelo Grok
def generar_imagen_grok(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (Grok)")
        
        # Importar la función original
        from modelos.IA_grok2_image import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))
    
# Función para generar una imagen con el modelo Gemini Flash
def generar_imagen_gemini_flash(section, content, descripcion, signals=None, print_log=True):
    # Función interna para manejar logs
    def log_message(msg):
        print(msg)
        if signals:
            signals.update_log.emit(str(msg))
    
    # Obtener el idioma actual
    current_language = 'es'
    if signals and hasattr(signals, 'current_language'):
        current_language = signals.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'current_language'):
        current_language = signals.parent.current_language
    elif signals and hasattr(signals, 'parent') and hasattr(signals.parent, 'parent') and hasattr(signals.parent.parent, 'current_language'):
        current_language = signals.parent.parent.current_language

    try:
        # Imprimir mensaje solo si se solicita
        if print_log:
            log_message(f"{obtener_traduccion('intentando_generar_imagen', current_language)} (GeminiFlash)")
        
        # Importar la función original
        from modelos.IA_gemini2_flash_image import generar_imagen
        return generar_imagen(section, content, descripcion)
    except Exception as e:
        raise RuntimeError(obtener_traduccion('error_generar_imagen', current_language).format(error=str(e)))